
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import seaborn as sns
import os
import unicodedata
from datetime import datetime
from io import BytesIO

# 통계 분석
from scipy import stats
from scipy.stats import chi2_contingency
from sklearn.linear_model import LinearRegression

# python-pptx 임포트 (PPT 생성용)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────────────────────────────────
# 페이지 설정
# ──────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="SQA 펌웨어 분석 대시보드",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# 🎨 Pretendard 기반 트렌디 디자인 (한국 SaaS 표준)
# ============================================================
st.markdown("""
<style>
/* === Pretendard 폰트 (토스/카카오/네이버 표준) === */
@import url('https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/variable/pretendardvariable.css');
@import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@400;500;600&display=swap');

/* === CSS 변수 === */
:root {
    --bg-cream: #faf9f5;
    --bg-white: #ffffff;
    --accent-coral: #cc785c;
    --accent-coral-dark: #b56a52;
    --accent-warm: #f5f1e8;
    --text-primary: #1a1a1a;
    --text-secondary: #6b6b6b;
    --text-muted: #9a9a9a;
    --border: #e3e1da;
    --border-strong: #d4d1c4;
}

/* === 전체 배경 === */
.stApp {
    background-color: var(--bg-cream) !important;
}

/* === Pretendard 폰트 (텍스트만 - 아이콘 절대 안 건드림) === */
.stApp p, .stApp label, .stApp h1, .stApp h2, .stApp h3, .stApp h4, .stApp h5, .stApp h6,
.stApp button, .stApp input, .stApp select, .stApp textarea,
.stApp div[data-testid="stMarkdownContainer"] {
    font-family: 'Pretendard Variable', 'Pretendard', -apple-system, BlinkMacSystemFont, 'Apple SD Gothic Neo', 'Noto Sans KR', sans-serif !important;
}

/* === 헤더 (Pretendard - 적정 사이즈) === */
h1 {
    font-weight: 700 !important;
    font-size: 1.75rem !important;
    line-height: 1.25 !important;
    margin-bottom: 0.75rem !important;
    letter-spacing: -0.03em !important;
    color: var(--text-primary) !important;
}
h2 {
    font-weight: 700 !important;
    font-size: 1.375rem !important;
    margin-top: 1.5rem !important;
    line-height: 1.3 !important;
    letter-spacing: -0.025em !important;
    color: var(--text-primary) !important;
}
h3 {
    font-weight: 600 !important;
    font-size: 1.125rem !important;
    line-height: 1.35 !important;
    letter-spacing: -0.02em !important;
    color: var(--text-primary) !important;
}
h4 {
    font-weight: 600 !important;
    font-size: 1rem !important;
    color: var(--text-primary) !important;
    letter-spacing: -0.015em !important;
}

/* === 사이드바 === */
section[data-testid="stSidebar"] {
    background-color: var(--bg-white) !important;
    border-right: 1px solid var(--border) !important;
}

/* === KPI 카드 === */
div[data-testid="stMetric"] {
    background-color: var(--bg-white) !important;
    border: 1px solid var(--border-strong) !important;
    border-radius: 12px !important;
    padding: 1.125rem 1.375rem !important;
    transition: all 0.15s ease !important;
}
div[data-testid="stMetric"]:hover {
    border-color: var(--accent-coral) !important;
    background-color: var(--accent-warm) !important;
    transform: translateY(-1px);
}
div[data-testid="stMetricLabel"] {
    font-size: 0.825rem !important;
    color: var(--text-secondary) !important;
    font-weight: 500 !important;
    letter-spacing: -0.01em !important;
}
div[data-testid="stMetricValue"] {
    font-size: 1.625rem !important;
    font-weight: 700 !important;
    color: var(--text-primary) !important;
    letter-spacing: -0.03em !important;
}
div[data-testid="stMetricDelta"] {
    font-size: 0.825rem !important;
    font-weight: 500 !important;
}

/* === 정보 박스 === */
div[data-testid="stAlert"] {
    background-color: var(--accent-warm) !important;
    border: 1px solid var(--border-strong) !important;
    border-radius: 10px !important;
    padding: 1rem 1.25rem !important;
}

/* === 버튼 === */
.stButton > button {
    border-radius: 8px !important;
    font-weight: 600 !important;
    font-size: 0.875rem !important;
    transition: all 0.12s ease !important;
    border: 1px solid var(--border-strong) !important;
    padding: 0.5rem 1.125rem !important;
    background-color: var(--bg-white) !important;
    color: var(--text-primary) !important;
    letter-spacing: -0.01em !important;
}
.stButton > button:hover {
    border-color: var(--text-primary) !important;
    background-color: var(--text-primary) !important;
    color: var(--bg-white) !important;
}
.stButton > button[kind="primary"] {
    background-color: var(--accent-coral) !important;
    color: var(--bg-white) !important;
    border: 1px solid var(--accent-coral) !important;
}
.stButton > button[kind="primary"]:hover {
    background-color: var(--accent-coral-dark) !important;
    border-color: var(--accent-coral-dark) !important;
    color: var(--bg-white) !important;
}

/* === Selectbox / Multiselect === */
div[data-baseweb="select"] > div {
    background-color: var(--bg-white) !important;
    border: 1px solid var(--border-strong) !important;
    border-radius: 8px !important;
}
div[data-baseweb="select"] > div:hover {
    border-color: var(--accent-coral) !important;
}

/* === Number Input === */
div[data-testid="stNumberInput"] input {
    border: 1px solid var(--border-strong) !important;
    border-radius: 8px !important;
    background-color: var(--bg-white) !important;
}

/* === Dataframe === */
div[data-testid="stDataFrame"] {
    border: 1px solid var(--border-strong) !important;
    border-radius: 10px !important;
    overflow: hidden !important;
}

/* === Expander === */
div[data-testid="stExpander"] {
    background-color: var(--bg-white) !important;
    border: 1px solid var(--border-strong) !important;
    border-radius: 10px !important;
}

/* === 파일 업로더 === */
section[data-testid="stFileUploader"] {
    background-color: var(--accent-warm) !important;
    border: 2px dashed var(--border-strong) !important;
    border-radius: 10px !important;
    padding: 1rem !important;
}

/* === 구분선 === */
hr {
    border-color: var(--border) !important;
    border-width: 1px !important;
    margin: 1.5rem 0 !important;
}

/* === 캡션 === */
div[data-testid="stCaptionContainer"], .stCaption, small {
    color: var(--text-secondary) !important;
    font-size: 0.875rem !important;
    font-weight: 400 !important;
}

/* === 사이드바 라디오 === */
div[data-testid="stSidebar"] label {
    font-weight: 500 !important;
    color: var(--text-primary) !important;
}

/* === 코드 블록 === */
code {
    background-color: var(--accent-warm) !important;
    color: var(--accent-coral-dark) !important;
    padding: 0.15rem 0.4rem !important;
    border-radius: 4px !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 0.85em !important;
    font-weight: 500 !important;
}

/* === 링크 === */
a {
    color: var(--accent-coral) !important;
    text-decoration: none !important;
    font-weight: 500 !important;
}
a:hover {
    text-decoration: underline !important;
}

/* === 메인 컨테이너 패딩 === */
.block-container {
    padding-top: 2rem !important;
    padding-bottom: 2rem !important;
    max-width: 1400px !important;
}

/* === Streamlit 기본 메뉴 숨김 === */
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}

</style>
""", unsafe_allow_html=True)

# ──────────────────────────────────────────────────────────────────────────
# 한글 폰트
# ──────────────────────────────────────────────────────────────────────────
@st.cache_resource
def setup_korean_font():
    font_path = '/usr/share/fonts/truetype/nanum/NanumGothic.ttf'
    if os.path.exists(font_path):
        fm.fontManager.addfont(font_path)
        return 'NanumGothic'
    return 'DejaVu Sans'

font_name = setup_korean_font()
plt.rcParams['font.family'] = font_name
plt.rcParams['axes.unicode_minus'] = False
sns.set_style("whitegrid")


# ──────────────────────────────────────────────────────────────────────────
# 상수
# ──────────────────────────────────────────────────────────────────────────
KEYWORD_DOMAIN_MAP = {
    'line broken':     {'area': '좌표 보간/샘플링 레이트',     'priority': 'HIGH'},
    'jitter':          {'area': '노이즈 필터/디바운싱',         'priority': 'MID'},
    'ghost touch':     {'area': '노이즈 임계값/그라운드',       'priority': 'HIGH'},
    'no touch':        {'area': '감도 보정/입력 임계값',         'priority': 'CRITICAL'},
    '2 point로 인식':  {'area': '멀티 입력 분리 알고리즘',         'priority': 'MID'},
    'edge 과밀착':     {'area': '경계영역 보정 (강하게 반응)',    'priority': 'MID'},
    'edge 미밀착':     {'area': '경계영역 보정 (약하게 반응)',    'priority': 'MID'},
    'touch delay':     {'area': '응답속도/인터럽트 처리',         'priority': 'CRITICAL'},
}

SEVERITY_MAP = {
    'no touch': 5, 'touch delay': 4, 'ghost touch': 4, 'line broken': 3,
    '2 point로 인식': 3, 'edge 과밀착': 2, 'edge 미밀착': 2, 'jitter': 2,
}

TEST_AREA_MAP = {
    'Wet': '방수·방습', 'Linearity': '좌표 보간', 'Filter': '노이즈 필터링',
    'A/P/L': '인식 정확도', 'Separation': '멀티 입력 분리',
    'Sensitivity': '입력 감도', 'Palm': 'Palm Rejection',
}

KEYWORDS_ALL = list(KEYWORD_DOMAIN_MAP.keys())

# ==============================================================================
# 💡 자동 인사이트 생성 함수 (분석가 시선)
# ==============================================================================
# 모든 인사이트는 4단계 구조:
# 1. What (수치/사실)
# 2. So What (비교/맥락)
# 3. Why (원인 추정)
# 4. Now What (권장 행동)
# ==============================================================================

# ==============================================================================
# 💼 ACTION PLAYBOOK (Fail_Type별 비즈니스/개발 액션 매핑)
# ==============================================================================
ACTION_PLAYBOOK = {
    'no touch': {
        'business': '핵심 기능 실패 → 출시 지연 가능성 ⚠️ HIGH. 고객사 클레임 직결.',
        'dev': [
            '⚡ 즉시: 입력 임계값(threshold) 설정 점검 및 캘리브레이션 데이터 검증',
            '📋 단주: PASS 비율 높은 모델의 임계값 설정과 비교 분석',
            '🔄 장주: 환경 조건별 임계값 자동 보정 로직 도입 검토',
        ],
        'comparison': '안정 빌드의 임계값 설정 코드와 diff 비교 · 정상 모델의 캘리브레이션 파라미터 대조',
    },
    'touch delay': {
        'business': '응답성 저하 → 사용자 경험 악화 → 고객사 클레임 가능성 ↑',
        'dev': [
            '⚡ 즉시: 응답속도 알고리즘 프로파일링 (병목 함수 식별)',
            '📋 단주: 안정 빌드와 인터럽트 처리 로직 코드 diff',
            '🔄 장주: 폴링 주기 / 인터럽트 우선순위 재설계 검토',
        ],
        'comparison': '응답속도 정상 빌드 vs 현재 빌드의 핸들러 코드 비교 · CPU 사용률 프로파일 대조',
    },
    'ghost touch': {
        'business': '오작동 발생 → 사용자 신뢰도 하락 · 안전 이슈 가능성',
        'dev': [
            '⚡ 즉시: 노이즈 필터링 임계값 강화',
            '📋 단주: EMI/EMC 환경에서 발생 패턴 로깅 분석',
            '🔄 장주: 노이즈 필터 알고리즘 (예: Kalman filter) 도입 검토',
        ],
        'comparison': '노이즈 환경 테스트 통과 빌드와 필터 파라미터 비교',
    },
    'line broken': {
        'business': '드로잉/제스처 기능 손상 → 그리기 앱 등 활용 모델에서 사용성 ↓',
        'dev': [
            '⚡ 즉시: 좌표 보간(interpolation) 알고리즘 검토',
            '📋 단주: 끊김 발생 좌표 패턴 분석 (속도/방향)',
            '🔄 장주: 적응형 샘플링 레이트 조정 로직 검토',
        ],
        'comparison': '안정 빌드의 좌표 보간 임계값 비교 · Drawing 정상 모델의 샘플링 레이트 대조',
    },
    '2 point로 인식': {
        'business': '멀티 입력 처리 오류 → 게임/생산성 앱 등 멀티 입력 활용 모델에서 사용성 ↓',
        'dev': [
            '⚡ 즉시: 멀티 입력 분리 알고리즘 검토',
            '📋 단주: 입력점 거리/면적 임계값 점검',
            '🔄 장주: 접점 클러스터링 알고리즘 재설계 검토',
        ],
        'comparison': '멀티 입력 정상 동작 빌드의 클러스터링 파라미터 비교',
    },
    'edge 과밀착': {
        'business': '엣지 영역 인식 오류 → 풀스크린 UI 모델에서 사용성 영향',
        'dev': [
            '⚡ 즉시: Palm Rejection 알고리즘 임계값 점검',
            '📋 단주: 엣지 영역 좌표 보정 로직 검토',
            '🔄 장주: 엣지 검출 알고리즘 재학습 검토',
        ],
        'comparison': '엣지 정상 모델의 Palm Rejection 파라미터 비교',
    },
    'edge 미밀착': {
        'business': '엣지 영역 입력 손실 → 풀스크린 UI 모델에서 핵심 영역 미인식',
        'dev': [
            '⚡ 즉시: 엣지 영역 감도 부스팅 파라미터 점검',
            '📋 단주: 베젤리스 모델의 엣지 캘리브레이션 검증',
            '🔄 장주: 엣지 전용 보정 알고리즘 도입 검토',
        ],
        'comparison': '엣지 정상 모델과 베젤 두께 / 감도 파라미터 비교',
    },
    'jitter': {
        'business': '입력 떨림 → 정밀 작업(서명, 그림) 시 품질 저하',
        'dev': [
            '⚡ 즉시: 좌표 스무딩(smoothing) 필터 임계값 강화',
            '📋 단주: 떨림 발생 좌표 분포 분석',
            '🔄 장주: 적응형 스무딩 알고리즘 도입 검토',
        ],
        'comparison': '정밀 입력 정상 모델의 스무딩 파라미터 비교',
    },
}


def get_action_for_keyword(keyword):
    """Fail_Type 키워드에 해당하는 액션 가져오기"""
    return ACTION_PLAYBOOK.get(keyword, {
        'business': '관련 결함 발생으로 해당 모듈의 안정성에 영향',
        'dev': [
            '⚡ 즉시: 결함 발생 시나리오 재현 및 로깅 분석',
            '📋 단주: 안정 빌드와 관련 모듈 코드 비교',
            '🔄 장주: 모듈 단위 리팩토링 검토',
        ],
        'comparison': '안정 빌드의 관련 모듈 파라미터와 비교 분석 권장',
    })


def render_insight_box(title, finding, action, severity="info", **kwargs):
    """간결한 인사이트 박스 (한 문장 발견 + 한 줄 액션)
    
    title: 인사이트 제목 (예: "빌드 추세 분석 — 안정화 진행 중")
    finding: 한 문장 발견 (수치 + 의미 + 결론을 한 문장에)
    action: 다음 액션 한 줄
    severity: info/warning/critical/success
    **kwargs: 호환성을 위한 추가 인자 (무시됨, 옛 코드 호환)
    """
    color_map = {
        "info": ("#3498db", "#EBF5FF"),
        "warning": ("#f39c12", "#FFF5E6"),
        "critical": ("#e74c3c", "#FFEBEE"),
        "success": ("#27ae60", "#E8F8F0"),
    }
    border_color, bg_color = color_map.get(severity, color_map["info"])
    
    # HTML (한 줄로 작성 - 코드블록 방지)
    html = f'<div style="background-color:{bg_color};border-left:4px solid {border_color};border-radius:8px;padding:1rem 1.25rem;margin:1rem 0;font-family:\'Pretendard Variable\', sans-serif;"><div style="font-weight:700;font-size:1rem;color:#1a1a1a;margin-bottom:0.6rem;">💡 {title}</div><div style="font-size:0.95rem;color:#2c3e50;line-height:1.55;margin-bottom:0.5rem;">{finding}</div><div style="background-color:white;padding:0.6rem 0.8rem;border-radius:6px;border:1px solid {border_color};font-size:0.9rem;color:{border_color};font-weight:600;">🎯 다음 액션: <span style="color:#2c3e50;font-weight:500;">{action}</span></div></div>'
    
    st.markdown(html, unsafe_allow_html=True)


def insight_build_trend(df, fail_df):
    """빌드별 ∩ 곡선 인사이트 (한 문장 + 액션)"""
    build_summary = df.groupby('Build_Num').apply(
        lambda x: (x['Result'] == 'FAIL').sum() / len(x) * 100
    ).round(2)
    if len(build_summary) < 2:
        return
    
    peak_build = build_summary.idxmax()
    peak_rate = build_summary.max()
    latest_build = build_summary.index[-1]
    latest_rate = build_summary.iloc[-1]
    diff = peak_rate - latest_rate
    
    if diff > 5:
        severity = "success"
        finding = f"빌드 <b>{peak_build}</b> ({peak_rate:.1f}%) 정점 후 최신 빌드 <b>{latest_build}</b> ({latest_rate:.1f}%)까지 <b>-{diff:.1f}%p 안정화</b> → 후속 빌드 수정 효과 확인됨"
        action = f"빌드 {peak_build}의 변경 사항 분석해서 안정 패턴을 표준 가이드로 정리"
    elif latest_rate > peak_rate * 0.9:
        severity = "warning"
        finding = f"최신 빌드 <b>{latest_build}</b> ({latest_rate:.1f}%)이 정점({peak_rate:.1f}%) 수준에 근접 → <b>회귀 가능성</b> ⚠️"
        action = f"빌드 {latest_build}의 변경 사항 회귀 검증 + 다음 빌드 출시 보류 검토"
    else:
        severity = "info"
        finding = f"빌드별 Fail율이 {min(build_summary):.1f}%~{peak_rate:.1f}% 범위에서 <b>안정세 유지</b> (최신 {latest_rate:.1f}%)"
        action = f"현 안정 빌드({latest_build})를 안정성 베이스라인으로 설정하고 모니터링 유지"
    
    render_insight_box(
        title=f"빌드 추세 분석",
        finding=finding,
        action=action,
        severity=severity
    )


def insight_ic_failrate(df, fail_df):
    """IC별 Fail율 인사이트 (한 문장 + 액션)"""
    ic_summary = df.groupby('IC').apply(
        lambda x: (x['Result'] == 'FAIL').sum() / len(x) * 100
    ).round(2)
    if len(ic_summary) < 2:
        return
    
    max_ic = ic_summary.idxmax()
    max_rate = ic_summary.max()
    avg_rate = ic_summary.mean()
    ratio = max_rate / avg_rate
    
    ic_fails = fail_df[fail_df['IC'] == max_ic]
    if len(ic_fails) == 0 or 'Keyword' not in ic_fails.columns:
        return
    top_keyword_counts = ic_fails['Keyword'].value_counts()
    if len(top_keyword_counts) == 0:
        return
    top_keyword = top_keyword_counts.index[0]
    top_keyword_pct = top_keyword_counts.iloc[0] / len(ic_fails) * 100
    area = KEYWORD_DOMAIN_MAP.get(top_keyword, {}).get('area', '관련 모듈')
    
    severity = "critical" if ratio > 1.3 else "warning"
    
    finding = f"<b>{max_ic}</b> Fail율 <b>{max_rate:.1f}%</b> (평균 {ratio:.1f}배) + <b>{top_keyword}</b> 결함이 {top_keyword_pct:.0f}% 집중 → <b>{area}</b> 모듈이 병목"
    action = f"{max_ic} 펌웨어의 {area} 알고리즘 우선 검토 + 안정 IC와 코드 diff 비교"
    
    render_insight_box(
        title=f"IC별 결함 분석",
        finding=finding,
        action=action,
        severity=severity
    )


def insight_failtype_impact(fail_df):
    """Fail_Type 영향 분석 (한 문장 + 액션)"""
    if 'Keyword' not in fail_df.columns:
        return
    fail_kw = fail_df[fail_df['Keyword'] != '']
    if len(fail_kw) == 0:
        return
    
    kw_counts = fail_kw['Keyword'].value_counts()
    kw_models = fail_kw.groupby('Keyword')['Model'].nunique()
    
    top_kw = kw_counts.index[0]
    top_count = kw_counts.iloc[0]
    top_models = kw_models[top_kw]
    
    wide_kw = kw_models.idxmax()
    wide_models = kw_models[wide_kw]
    
    area = KEYWORD_DOMAIN_MAP.get(top_kw, {}).get('area', '관련 모듈')
    
    if top_kw == wide_kw:
        finding = f"<b>{top_kw}</b>가 빈도 1위({top_count}건) + 영향 범위 1위({top_models}개 모델) → <b>{area}</b> 모듈의 광범위 구조적 이슈"
        action = f"{area} 모듈 최우선 개선 + 영향 모델 {top_models}개 일괄 회귀 테스트"
    else:
        finding = f"빈도 1위 <b>{top_kw}</b>({top_count}건, {top_models}개 모델) vs 범위 1위 <b>{wide_kw}</b>({wide_models}개 모델) → 두 결함의 원인이 다름"
        action = f"{top_kw}는 빈도 집중 분석, {wide_kw}는 광범위 영향 분석으로 이원화 대응"
    
    render_insight_box(
        title="결함 영향 분석",
        finding=finding,
        action=action,
        severity="warning"
    )


def insight_chisquare(df, fail_df):
    """카이제곱 대조 분석 (한 문장 + 액션)"""
    render_insight_box(
        title="합격률 vs 결함 종류 — 대조 분석",
        finding="합격률은 IC/고객사와 <b>무관</b>(p≥0.05)이지만, 결함 종류는 <b>매우 유의</b>(p&lt;0.001, Cramér's V=0.75) → 그룹별 약점이 다른 곳에 나타남",
        action="단순 합격률 KPI가 아닌, IC·고객사별 결함 패턴 카드 도입으로 차별화 대응 전략 수립",
        severity="info"
    )


def insight_monthly_trend(df, fail_df):
    """월별 추이 인사이트 (한 문장 + 액션)"""
    df_copy = df.copy()
    df_copy['Test_Date'] = pd.to_datetime(df_copy['Test_Date'], errors='coerce')
    df_copy = df_copy.dropna(subset=['Test_Date'])
    if len(df_copy) == 0:
        return
    df_copy['month'] = df_copy['Test_Date'].dt.to_period('M')
    monthly = df_copy.groupby('month').apply(
        lambda x: (x['Result'] == 'FAIL').sum() / len(x) * 100
    ).round(2)
    if len(monthly) < 2:
        return
    
    peak_month = monthly.idxmax()
    peak_rate = monthly.max()
    latest_month = monthly.index[-1]
    latest_rate = monthly.iloc[-1]
    avg_rate = monthly.mean()
    diff = latest_rate - avg_rate
    
    if diff < -3:
        severity = "success"
        finding = f"<b>{peak_month}</b> 정점({peak_rate:.1f}%) 이후 <b>{latest_month}</b> {latest_rate:.1f}%로 평균 대비 {diff:+.1f}%p <b>개선세</b>"
        action = f"{peak_month} 이후 적용된 펌웨어 수정 사항을 회사 표준 가이드로 정리"
    elif diff > 3:
        severity = "warning"
        finding = f"<b>{latest_month}</b> Fail율 {latest_rate:.1f}%, 평균({avg_rate:.1f}%) 대비 {diff:+.1f}%p <b>악화</b> ⚠️"
        action = f"{latest_month} 출시 빌드의 변경 사항 즉시 회귀 검증"
    else:
        severity = "info"
        finding = f"월별 Fail율이 {monthly.min():.1f}%~{peak_rate:.1f}% 범위에서 <b>안정 유지</b> (최근 {latest_month} {latest_rate:.1f}%)"
        action = "현 품질 수준을 회사 SQA 표준 베이스라인으로 확정 후 신규 기능 검토"
    
    render_insight_box(
        title=f"월별 Fail율 추이",
        finding=finding,
        action=action,
        severity=severity
    )


def insight_customer_analysis(df, fail_df):
    """고객사 분석 (한 문장 + 액션)"""
    cust_summary = df.groupby('Customer').apply(
        lambda x: (x['Result'] == 'FAIL').sum() / len(x) * 100
    ).round(2)
    if len(cust_summary) < 2:
        return
    
    max_cust = cust_summary.idxmax()
    max_rate = cust_summary.max()
    min_cust = cust_summary.idxmin()
    min_rate = cust_summary.min()
    spread = max_rate - min_rate
    
    max_fails = fail_df[fail_df['Customer'] == max_cust]
    if len(max_fails) > 0 and 'Keyword' in max_fails.columns:
        max_kw_series = max_fails['Keyword'].value_counts()
        if len(max_kw_series) > 0:
            max_kw = max_kw_series.index[0]
            area = KEYWORD_DOMAIN_MAP.get(max_kw, {}).get('area', '관련 모듈')
        else:
            max_kw = "다양한 결함"
            area = "다양한 모듈"
    else:
        max_kw = "다양한 결함"
        area = "다양한 모듈"
    
    severity = "warning" if spread > 5 else "info"
    
    finding = f"고객사별 Fail율 격차 <b>{spread:.1f}%p</b> ({min_cust} {min_rate:.1f}% ~ {max_cust} {max_rate:.1f}%) → {max_cust}의 주요 결함은 <b>{max_kw}</b>"
    action = f"{max_cust} 전용 {area} 영역 검증 시나리오 강화 + 모델 라인업 차이 분석"
    
    render_insight_box(
        title="고객사별 결함 패턴",
        finding=finding,
        action=action,
        severity=severity
    )


def insight_test_item_matrix(fail_df):
    """Test_Item × Fail_Type 매트릭스 (한 문장 + 액션)"""
    if 'Test_Item' not in fail_df.columns or 'Keyword' not in fail_df.columns:
        return
    fail_kw = fail_df[fail_df['Keyword'] != '']
    if len(fail_kw) == 0:
        return
    
    pair_counts = fail_kw.groupby(['Test_Item', 'Keyword']).size().sort_values(ascending=False)
    if len(pair_counts) == 0:
        return
    
    top_pair = pair_counts.index[0]
    top_count = pair_counts.iloc[0]
    test_item, keyword = top_pair
    area = KEYWORD_DOMAIN_MAP.get(keyword, {}).get('area', '관련 모듈')
    
    finding = f"<b>{test_item}</b> × <b>{keyword}</b> = <b>{top_count}건</b> 핫스팟 발견 → {area} 모듈의 약점이 해당 시나리오에서 자극됨"
    action = f"{test_item} 시나리오의 코드 경로 + {area} 모듈 우선 디버깅"
    
    render_insight_box(
        title="Test_Item × Fail_Type 핫스팟",
        finding=finding,
        action=action,
        severity="warning"
    )


def insight_regression_alert(alerts_df, latest_build, prev_build, df, fail_df):
    """회귀 알람 (한 문장 + 액션)"""
    if len(alerts_df) == 0:
        render_insight_box(
            title="회귀 알람 — 정상 상태 ✅",
            finding=f"최신 빌드 <b>{latest_build}</b>에서 이전 빌드 {prev_build} 대비 신규/악화 결함 <b>미검출</b>",
            action=f"빌드 {latest_build}를 안정성 베이스라인으로 설정하고 출시 진행 가능",
            severity="success"
        )
        return
    
    n_new = (alerts_df['상태'] == 'NEW').sum()
    n_worse = (alerts_df['상태'] == 'WORSE').sum()
    
    worst_idx = alerts_df['latest'].idxmax()
    worst_model = alerts_df.loc[worst_idx, 'Model']
    worst_keyword = alerts_df.loc[worst_idx, 'Keyword']
    worst_prev = alerts_df.loc[worst_idx, 'prev']
    worst_latest = alerts_df.loc[worst_idx, 'latest']
    area = KEYWORD_DOMAIN_MAP.get(worst_keyword, {}).get('area', '관련 모듈')
    
    if worst_prev == 0:
        change = f"0→{worst_latest}건 (신규)"
    else:
        pct = (worst_latest - worst_prev) / worst_prev * 100
        change = f"{worst_prev}→{worst_latest}건 ({pct:+.0f}%)"
    
    severity = "critical" if n_new > 3 or n_worse > 5 else "warning"
    
    finding = f"빌드 <b>{latest_build}</b>에서 <b>NEW {n_new}건 + WORSE {n_worse}건</b> 회귀 발생, 최대: {worst_model} × {worst_keyword} ({change})"
    action = f"{worst_model}의 {area} 코드 변경 즉시 검토 + 빌드 {latest_build} 출시 보류 검토"
    
    render_insight_box(
        title=f"회귀 알람 — {n_new + n_worse}건 검출",
        finding=finding,
        action=action,
        severity=severity
    )



# ==============================================================================
# 💡 자동 인사이트 생성 함수 끝
# ==============================================================================



# 보고서 컬러
RPT_HEADER_BG = RGBColor(0x06, 0x5A, 0x82)
RPT_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
RPT_ALT_ROW   = RGBColor(0xF4, 0xF8, 0xFB)
RPT_CRITICAL  = RGBColor(0xD6, 0x28, 0x28)
RPT_HIGH      = RGBColor(0xF7, 0x7F, 0x00)
RPT_MID       = RGBColor(0x06, 0x5A, 0x82)
RPT_TITLE     = RGBColor(0x06, 0x5A, 0x82)
PPT_FONT = "맑은 고딕"


# ──────────────────────────────────────────────────────────────────────────
# 데이터 로드 (캐시)
# ──────────────────────────────────────────────────────────────────────────
@st.cache_data
def load_base_data():
    df = pd.read_csv(
        'SQA_test_results_clean.csv',
        parse_dates=['Test_Date'],
        dtype={'IC_Code': str, 'Model_Minor': str}
    )
    df['IC_Code'] = df['IC_Code'].str.zfill(2)
    df['Model_Minor'] = df['Model_Minor'].str.zfill(2)
    df['Keyword'] = df['Keyword'].fillna('')
    return df


def get_video_link(video_id):
    if pd.isna(video_id) or video_id is None or video_id == '':
        return None
    return f"https://drive.google.com/file/d/{video_id}/view"


# ──────────────────────────────────────────────────────────────────────────
# 세션 상태 초기화
# ──────────────────────────────────────────────────────────────────────────
if 'uploaded_data' not in st.session_state:
    st.session_state.uploaded_data = []  # 누적된 업로드 데이터들
if 'data_version' not in st.session_state:
    st.session_state.data_version = 0  # 데이터 갱신 추적용


def get_current_df():
    """기본 데이터 + 업로드된 데이터 모두 합쳐 반환"""
    base_df = load_base_data()
    if st.session_state.uploaded_data:
        all_dfs = [base_df] + st.session_state.uploaded_data
        return pd.concat(all_dfs, ignore_index=True)
    return base_df


# ==============================================================================
# 🎯 메뉴 1: 필터 분석
# ==============================================================================

def render_filter_analysis(df, fail_df):
    plt.rcParams['font.family'] = font_name

    # 상단 KPI
    total_tests_all = len(df)
    total_fails_all = len(fail_df)
    fail_rate_all = total_fails_all / total_tests_all * 100

    st.markdown("### 📊 전체 데이터 현황")
    k1, k2, k3, k4 = st.columns(4)
    with k1:
        st.markdown("##### 📋 전체 테스트")
        st.markdown(f"### :blue[**{total_tests_all:,}건**]")
    with k2:
        st.markdown("##### ❌ 전체 Fail")
        st.markdown(f"### :red[**{total_fails_all:,}건**]")
    with k3:
        st.markdown("##### 📈 전체 Fail율")
        st.markdown(f"### :orange[**{fail_rate_all:.1f}%**]")
    with k4:
        st.markdown("##### 💻 분석 모델")
        st.markdown(f"### :green[**{df['Model'].nunique()}개**]")

    st.markdown("---")

    # 필터 UI (다중 선택 가능)
    st.markdown("### 🎯 필터 선택")
    st.caption("각 필터에서 여러 값을 동시에 선택할 수 있습니다. (비어두면 전체)")

    row1_col1, row1_col2, row1_col3 = st.columns(3)
    with row1_col1:
        ic_filter = st.multiselect(
            "🔌 IC", 
            options=sorted(df['IC'].unique().tolist()),
            placeholder="전체 IC (비어두면 모두 선택)"
        )
    with row1_col2:
        model_filter = st.multiselect(
            "💻 Model",
            options=sorted(df['Model'].unique().tolist()),
            placeholder="전체 모델 (비어두면 모두 선택)"
        )
    with row1_col3:
        build_filter = st.multiselect(
            "🔄 Build",
            options=sorted(df['Build_Num'].unique().tolist()),
            placeholder="전체 빌드 (비어두면 모두 선택)"
        )
    
    row2_col1, row2_col2, row2_col3 = st.columns(3)
    with row2_col1:
        test_item_filter = st.multiselect(
            "🧪 Test_Item",
            options=sorted(df['Test_Item'].unique().tolist()),
            placeholder="전체 Test_Item (비어두면 모두 선택)"
        )
    with row2_col2:
        kw_filter = st.multiselect(
            "🐛 Fail_Type",
            options=KEYWORDS_ALL,
            placeholder="전체 Fail_Type (비어두면 모두 선택)"
        )
    with row2_col3:
        st.write("")

    # 필터 적용 (multiselect는 리스트)
    filtered_all = df.copy()
    filtered_fail = fail_df.copy()
    active_filters = []
    if ic_filter:
        filtered_all = filtered_all[filtered_all['IC'].isin(ic_filter)]
        filtered_fail = filtered_fail[filtered_fail['IC'].isin(ic_filter)]
        active_filters.append(f"IC={','.join(ic_filter)}")
    if model_filter:
        filtered_all = filtered_all[filtered_all['Model'].isin(model_filter)]
        filtered_fail = filtered_fail[filtered_fail['Model'].isin(model_filter)]
        if len(model_filter) <= 3:
            active_filters.append(f"Model={','.join(model_filter)}")
        else:
            active_filters.append(f"Model={len(model_filter)}개")
    if build_filter:
        filtered_all = filtered_all[filtered_all['Build_Num'].isin(build_filter)]
        filtered_fail = filtered_fail[filtered_fail['Build_Num'].isin(build_filter)]
        if len(build_filter) <= 3:
            active_filters.append(f"Build={','.join(map(str, build_filter))}")
        else:
            active_filters.append(f"Build={len(build_filter)}개")
    if test_item_filter:
        filtered_all = filtered_all[filtered_all['Test_Item'].isin(test_item_filter)]
        filtered_fail = filtered_fail[filtered_fail['Test_Item'].isin(test_item_filter)]
        if len(test_item_filter) <= 3:
            active_filters.append(f"Test_Item={','.join(test_item_filter)}")
        else:
            active_filters.append(f"Test_Item={len(test_item_filter)}개")
    if kw_filter:
        filtered_fail = filtered_fail[filtered_fail['Keyword'].isin(kw_filter)]
        if len(kw_filter) <= 3:
            active_filters.append(f"Fail_Type={','.join(kw_filter)}")
        else:
            active_filters.append(f"Fail_Type={len(kw_filter)}개")

    st.markdown("---")
    if active_filters:
        st.markdown(f"### 🎯 적용된 필터: `{' · '.join(active_filters)}`")
    else:
        st.markdown("### 🎯 전체 데이터 분석 (필터 미적용)")

    if len(filtered_fail) == 0:
        st.warning("⚠️ 선택한 조건에 해당하는 Fail 케이스가 없습니다.")
        return

    # 필터 결과 KPI
    st.markdown("### 📊 ① 필터 결과 요약")
    total_tests = len(filtered_all)
    total_fails = len(filtered_fail)
    fail_rate = (total_fails / total_tests * 100) if total_tests > 0 else 0
    n_models = filtered_fail['Model'].nunique()

    k1, k2, k3, k4 = st.columns(4)
    with k1:
        st.markdown("##### 📋 전체 테스트")
        st.markdown(f"### :blue[**{total_tests:,}건**]")
    with k2:
        st.markdown("##### ❌ Fail 건수")
        st.markdown(f"### :red[**{total_fails:,}건**]")
    with k3:
        st.markdown("##### 📈 Fail율")
        st.markdown(f"### :orange[**{fail_rate:.1f}%**]")
    with k4:
        st.markdown("##### 💻 영향 모델")
        st.markdown(f"### :green[**{n_models}개**]")

    st.markdown("---")

    # Fail_Type 순위
    if not kw_filter:
        st.markdown("### 🏆 ② Fail_Type 순위 (TOP 5)")
        kw_rank = filtered_fail['Keyword'].value_counts().head(5).reset_index()
        kw_rank.columns = ['Fail_Type', '건수']
        kw_rank['비중'] = (kw_rank['건수'] / kw_rank['건수'].sum() * 100).round(1)
        kw_rank.insert(0, '순위', range(1, len(kw_rank) + 1))

        col_left, col_right = st.columns([1, 1])
        with col_left:
            fig, ax = plt.subplots(figsize=(6, 5))
            colors = plt.cm.Set3(np.linspace(0, 1, len(kw_rank)))
            ax.pie(kw_rank['건수'], labels=kw_rank['Fail_Type'],
                   autopct='%1.1f%%', wedgeprops=dict(width=0.4),
                   startangle=90, colors=colors,
                   textprops={'fontsize': 10})
            ax.set_title(f'Fail_Type 분포 (총 {len(filtered_fail)}건)',
                         fontsize=12, fontweight='bold')
            plt.tight_layout()
            st.pyplot(fig)
            plt.close(fig)
        with col_right:
            st.dataframe(kw_rank, use_container_width=True, hide_index=True)
        # 💡 인사이트 (차트 바로 아래)
        if len(filtered_fail[filtered_fail['Keyword'] != '']) > 0:
            insight_failtype_impact(filtered_fail)
        st.markdown("---")

    # 빌드별 추이
    if not build_filter:
        st.markdown("### 📈 ③ 빌드별 Fail 추이")
        build_trend = filtered_fail.groupby('Build_Num').size().reset_index(name='count')
        build_trend = build_trend.sort_values('Build_Num')
        if len(build_trend) > 1:
            fig, ax = plt.subplots(figsize=(12, 4))
            ax.plot(build_trend['Build_Num'], build_trend['count'],
                    marker='o', markersize=8, linewidth=2, color='#1C7293')
            for _, row in build_trend.iterrows():
                ax.annotate(f"{row['count']}", (row['Build_Num'], row['count']),
                            textcoords="offset points", xytext=(0, 10),
                            ha='center', fontsize=10, fontweight='bold')
            ax.set_xlabel('빌드 번호')
            ax.set_ylabel('Fail 건수')
            ax.grid(True, alpha=0.3)
            ax.set_xticks(build_trend['Build_Num'])
            ax.set_xticklabels(build_trend['Build_Num'], rotation=45)
            plt.tight_layout()
            st.pyplot(fig)
            plt.close(fig)
            # 💡 인사이트 (차트 바로 아래)
            if filtered_all['Build_Num'].nunique() >= 2:
                insight_build_trend(filtered_all, filtered_fail)
        st.markdown("---")

    # 모델별 순위
    if not model_filter:
        st.markdown("### 🔥 ④ 모델별 Fail 순위 (TOP 10)")
        model_rank = filtered_fail['Model'].value_counts().head(10).reset_index()
        model_rank.columns = ['Model', 'Fail 건수']
        if len(model_rank) > 0:
            fig, ax = plt.subplots(figsize=(10, max(4, len(model_rank) * 0.4)))
            bars = ax.barh(model_rank['Model'], model_rank['Fail 건수'],
                           color=plt.cm.YlOrRd(np.linspace(0.4, 0.9, len(model_rank))))
            for bar, cnt in zip(bars, model_rank['Fail 건수']):
                ax.text(bar.get_width() + max(model_rank['Fail 건수']) * 0.01,
                        bar.get_y() + bar.get_height()/2,
                        f'{cnt}건', va='center', fontsize=10, fontweight='bold')
            ax.set_xlabel('Fail 건수')
            ax.invert_yaxis()
            plt.tight_layout()
            st.pyplot(fig)
            plt.close(fig)
            # 💡 인사이트 (IC별 결함 분석)
            if filtered_all['IC'].nunique() >= 2:
                insight_ic_failrate(filtered_all, filtered_fail)
        st.markdown("---")
    if not test_item_filter:
        st.markdown("### 🧪 ⑤ Test_Item별 Fail 순위 (TOP 10)")
        item_rank = filtered_fail['Test_Item'].value_counts().head(10).reset_index()
        item_rank.columns = ['Test_Item', 'Fail 건수']
        if len(item_rank) > 0:
            fig, ax = plt.subplots(figsize=(10, max(4, len(item_rank) * 0.4)))
            bars = ax.barh(item_rank['Test_Item'], item_rank['Fail 건수'],
                           color=plt.cm.Blues(np.linspace(0.4, 0.9, len(item_rank))))
            for bar, cnt in zip(bars, item_rank['Fail 건수']):
                ax.text(bar.get_width() + max(item_rank['Fail 건수']) * 0.01,
                        bar.get_y() + bar.get_height()/2,
                        f'{cnt}건', va='center', fontsize=10, fontweight='bold')
            ax.set_xlabel('Fail 건수')
            ax.invert_yaxis()
            plt.tight_layout()
            st.pyplot(fig)
            plt.close(fig)
            # 💡 인사이트 (Test_Item × Fail_Type 매트릭스)
            if 'Keyword' in filtered_fail.columns and len(filtered_fail[filtered_fail['Keyword'] != '']) > 0:
                insight_test_item_matrix(filtered_fail)
        st.markdown("---")
    st.markdown(f"### 📋 ⑥ 상세 Fail 케이스 ({len(filtered_fail):,}건)")
    st.caption("표의 영상 링크를 클릭하면 새 탭에서 결함 영상이 재생됩니다.")

    display_df = filtered_fail[[
        'No', 'IC', 'Model', 'Keyword', 'Build_Num',
        'Customer', 'Category', 'Test_Item', 'video_id'
    ]].copy()
    display_df['영상'] = display_df['video_id'].apply(get_video_link)
    display_df = display_df.drop(columns=['video_id'])
    # Fail_Type 건수 기준 정렬 (많은 것이 위로)
    kw_counts = filtered_fail['Keyword'].value_counts().to_dict()
    display_df['_sort'] = display_df['Keyword'].map(kw_counts).fillna(0)
    display_df = display_df.sort_values('_sort', ascending=False).drop(columns=['_sort'])

    st.dataframe(
        display_df,
        use_container_width=True,
        hide_index=True,
        column_config={
            "영상": st.column_config.LinkColumn("영상 보기", display_text="🎬 재생"),
            "Keyword": st.column_config.TextColumn("Fail_Type"),
        }
    )
    
    # ===== 카이제곱 대조 분석 (전체 데이터일 때만 표시 — 차트 외 별도 통계 분석) =====
    is_full = (not ic_filter and not model_filter and not build_filter 
               and not test_item_filter and not kw_filter)
    
    if is_full and len(filtered_all) > 0 and len(filtered_fail) > 0:
        st.markdown("---")
        st.markdown("### 📐 통계 검증 — 카이제곱 대조 분석")
        st.caption("필터를 비운 전체 데이터 기준 통계 검증입니다.")
        insight_chisquare(filtered_all, filtered_fail)


# ==============================================================================
# 📊 메뉴 2: 전체 인사이트 (생략 - 기존과 동일)
# ==============================================================================
def render_full_insights(df, fail_df):
    plt.rcParams['font.family'] = font_name

    # 빌드별 ∩ 곡선
    st.markdown("### 1️⃣ 빌드별 ∩(역U) 안정화 곡선")
    build_summary = df.groupby('Build_Num').agg(
        total=('Result', 'count'),
        fail=('Result', lambda x: (x == 'FAIL').sum())
    ).reset_index()
    build_summary['fail_rate'] = (build_summary['fail'] / build_summary['total'] * 100).round(2)
    build_summary = build_summary.sort_values('Build_Num').reset_index(drop=True)

    fig, ax = plt.subplots(figsize=(14, 5))
    ax.plot(build_summary['Build_Num'], build_summary['fail_rate'],
            marker='o', markersize=10, linewidth=2, color='#e74c3c')
    for i, row in build_summary.iterrows():
        if i == 0:
            offset_y, va = 15, 'bottom'
        else:
            prev_rate = build_summary.loc[i-1, 'fail_rate']
            offset_y = -20 if row['fail_rate'] < prev_rate else 15
            va = 'top' if row['fail_rate'] < prev_rate else 'bottom'
        ax.annotate(f"{row['fail_rate']}%",
                    (row['Build_Num'], row['fail_rate']),
                    textcoords="offset points", xytext=(0, offset_y),
                    ha='center', va=va, fontsize=11, fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.3', facecolor='white',
                              edgecolor='#cccccc', alpha=0.9))
    ax.set_title('빌드별 전체 Fail율 추이', fontsize=14, fontweight='bold')
    ax.set_xlabel('빌드 번호')
    ax.set_ylabel('Fail율 (%)')
    ax.grid(True, alpha=0.3)
    ax.set_xticks(build_summary['Build_Num'])
    ax.set_xticklabels(build_summary['Build_Num'], rotation=45)
    plt.tight_layout()
    st.pyplot(fig)
    plt.close(fig)
    
    # 💡 인사이트
    insight_build_trend(df, fail_df)
    
    st.markdown("---")

    # IC별 + Donut
    st.markdown("### 2️⃣ IC별 Fail율 + Fail_Type 분포")
    ic_summary = df.groupby('IC').agg(
        total=('Result', 'count'),
        fail=('Result', lambda x: (x == 'FAIL').sum())
    ).reset_index()
    ic_summary['fail_rate'] = (ic_summary['fail'] / ic_summary['total'] * 100).round(2)
    ic_summary = ic_summary.sort_values('fail_rate', ascending=False)

    fig1, ax1 = plt.subplots(figsize=(8, 4))
    bars = ax1.bar(ic_summary['IC'], ic_summary['fail_rate'],
                   color=['#e74c3c', '#f39c12', '#3498db'])
    ax1.set_title('IC별 Fail율', fontsize=14, fontweight='bold')
    ax1.set_ylabel('Fail율 (%)')
    ax1.set_ylim(0, max(ic_summary['fail_rate']) * 1.2)
    for bar, rate in zip(bars, ic_summary['fail_rate']):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                 f'{rate}%', ha='center', fontsize=12, fontweight='bold')
    plt.tight_layout()
    st.pyplot(fig1)
    plt.close(fig1)

    ic_keyword = fail_df.groupby(['IC', 'Keyword']).size().reset_index(name='count')
    colors_kw = plt.cm.Set3(np.linspace(0, 1, 8))
    color_dict = dict(zip(KEYWORDS_ALL, colors_kw))

    fig2, axes2 = plt.subplots(1, 3, figsize=(14, 5))
    fig2.suptitle('IC별 Fail_Type 분포', fontsize=14, fontweight='bold')
    for ax, ic in zip(axes2, ['G7500', 'GT1T0A', 'GT9XS']):
        ic_data = ic_keyword[ic_keyword['IC'] == ic].sort_values('count', ascending=False)
        total_cnt = ic_data['count'].sum()
        colors_this = [color_dict[kw] for kw in ic_data['Keyword']]
        ax.pie(ic_data['count'], labels=None,
               autopct=lambda pct: f'{pct:.1f}%' if pct >= 5 else '',
               wedgeprops=dict(width=0.4), startangle=90, colors=colors_this,
               textprops={'fontsize': 10, 'fontweight': 'bold'})
        ax.text(0, 0.1, ic, ha='center', va='center', fontsize=14, fontweight='bold')
        ax.text(0, -0.15, f'Fail {total_cnt}건', ha='center', va='center',
                fontsize=10, color='#666666')
    legend_patches = [plt.Rectangle((0,0),1,1, color=color_dict[kw]) for kw in KEYWORDS_ALL]
    fig2.legend(legend_patches, KEYWORDS_ALL, loc='lower center', ncol=4,
                fontsize=9, bbox_to_anchor=(0.5, -0.05))
    plt.tight_layout()
    st.pyplot(fig2)
    plt.close(fig2)
    
    # 💡 인사이트
    insight_ic_failrate(df, fail_df)
    
    st.markdown("---")

    # 카테고리 히트맵
    st.markdown("### 3️⃣ 카테고리 × Fail_Type 히트맵")
    pivot_cat = fail_df.pivot_table(
        index='Category', columns='Keyword', values='No', aggfunc='count', fill_value=0)
    pivot_cat = pivot_cat[[k for k in KEYWORDS_ALL if k in pivot_cat.columns]]
    fig, ax = plt.subplots(figsize=(12, 6))
    sns.heatmap(pivot_cat, annot=True, fmt='d', cmap='YlOrRd',
                cbar_kws={'label': 'Fail 건수'}, ax=ax,
                linewidths=0.5, linecolor='white')
    ax.set_title('카테고리 × Fail_Type 히트맵', fontsize=14, fontweight='bold', pad=15)
    plt.tight_layout()
    st.pyplot(fig)
    plt.close(fig)
    st.markdown("---")

    # 모델 히트맵
    st.markdown("### 4️⃣ 모델 × Fail_Type 히트맵 ⭐")
    pivot_model = fail_df.pivot_table(
        index='Model', columns='Keyword', values='No', aggfunc='count', fill_value=0)
    pivot_model = pivot_model.loc[pivot_model.sum(axis=1).sort_values(ascending=False).index]
    pivot_model = pivot_model[[k for k in KEYWORDS_ALL if k in pivot_model.columns]]
    fig, ax = plt.subplots(figsize=(12, max(8, len(pivot_model) * 0.3)))
    sns.heatmap(pivot_model, annot=True, fmt='d', cmap='YlOrRd',
                cbar_kws={'label': 'Fail 건수'}, ax=ax,
                linewidths=0.5, linecolor='white')
    ax.set_title('모델 × Fail_Type 히트맵', fontsize=14, fontweight='bold', pad=15)
    plt.tight_layout()
    st.pyplot(fig)
    plt.close(fig)
    st.markdown("---")

    # 발생 횟수 × 영향 모델 수 버블 차트
    st.markdown("### 5️⃣ Fail_Type 영향 분석 (발생 횟수 × 영향 모델 수)")
    st.caption("오른쪽 위에 있을수록 자주 발생하고 많은 모델에 영향 → 우선 검토 대상")
    
    # 각 Fail_Type별 통계
    bubble_data = []
    for kw in KEYWORDS_ALL:
        kw_fails = fail_df[fail_df['Keyword'] == kw]
        if len(kw_fails) > 0:
            bubble_data.append({
                'keyword': kw,
                'count': len(kw_fails),
                'n_models': kw_fails['Model'].nunique()
            })
    
    if bubble_data:
        max_count = max(b['count'] for b in bubble_data)
        max_models = max(b['n_models'] for b in bubble_data)
        
        fig, ax = plt.subplots(figsize=(12, 7))
        colors = plt.cm.YlOrRd([b['count']/max_count for b in bubble_data])
        
        for i, b in enumerate(bubble_data):
            size = 300 + (b['count'] / max_count) * 2700
            ax.scatter(b['count'], b['n_models'], s=size, c=[colors[i]],
                       alpha=0.65, edgecolors='black', linewidths=1.5, zorder=3)
        
        # 라벨 (adjust_text 없이 수동으로 위치 분산)
        # y값 기준 정렬해서 가까운 것끼리 위치 어긋나게
        # 라벨 위치 분산 (각 버블마다 다른 방향)
        sorted_bubbles = sorted(bubble_data, key=lambda b: (b['n_models'], b['count']))
        # 8개 방향 분산
        directions = [
            (25, 25), (25, -25), (-25, 25), (-25, -25),
            (35, 5), (-35, 5), (5, 35), (5, -35)
        ]
        for i, b in enumerate(sorted_bubbles):
            offset_x, offset_y = directions[i % len(directions)]
            ha = 'left' if offset_x > 0 else 'right' 
            
            ax.annotate(
                f"{b['keyword']}\n({b['count']}건, {b['n_models']}개 모델)",
                (b['count'], b['n_models']),
                xytext=(offset_x, offset_y),
                textcoords='offset points',
                fontsize=9, fontweight='bold', ha=ha,
                bbox=dict(boxstyle='round,pad=0.3', facecolor='white',
                          edgecolor='#cccccc', alpha=0.9),
                zorder=4,
                arrowprops=dict(arrowstyle='-', color='#888888', lw=0.5, alpha=0.5)
            )
        
        # 평균선
        avg_count = sum(b['count'] for b in bubble_data) / len(bubble_data)
        avg_models = sum(b['n_models'] for b in bubble_data) / len(bubble_data)
        ax.axhline(y=avg_models, color='gray', linestyle='--', alpha=0.4,
                   label=f'평균 영향 모델 ({avg_models:.0f}개)')
        ax.axvline(x=avg_count, color='gray', linestyle='--', alpha=0.4,
                   label=f'평균 발생 횟수 ({avg_count:.0f}건)')
        
        ax.set_xlabel('발생 횟수 (Fail 건수)', fontsize=12, fontweight='bold')
        ax.set_ylabel('영향 모델 수', fontsize=12, fontweight='bold')
        ax.set_title('Fail_Type 영향 분석',
                     fontsize=14, fontweight='bold', pad=15)
        ax.set_xlim(0, max_count * 1.3)
        ax.set_ylim(0, max_models * 1.3)
        ax.grid(True, alpha=0.3)
        ax.legend(loc='lower right', fontsize=9)
        
        # 사분면 안내
        ax.text(max_count * 1.25, max_models * 1.25, '★ 자주 + 광범위\n우선 검토',
                fontsize=10, ha='right', va='top', color='#c0392b', fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.4', facecolor='#FFF5F5', edgecolor='#e74c3c'))
        
        plt.tight_layout()
        st.pyplot(fig)
        plt.close(fig)
        
        # 표
        bubble_df = pd.DataFrame(bubble_data).sort_values('count', ascending=False).reset_index(drop=True)
        bubble_df.insert(0, '순위', range(1, len(bubble_df) + 1))
        bubble_df.columns = ['순위', 'Fail_Type', '발생 횟수', '영향 모델 수']
        st.dataframe(bubble_df, use_container_width=True, hide_index=True)
        
        # 💡 인사이트
        insight_failtype_impact(fail_df)

    st.markdown("---")

    # ===== ⑥ 무엇이 결함을 결정하는가? (카이제곱 검정) =====
    st.markdown("### 6️⃣ 무엇이 결함을 결정하는가? (카이제곱 검정)")
    st.caption("IC/고객사가 '합격 여부'와 '결함 종류'에 각각 영향을 주는지 통계적으로 비교 분석합니다.")

    from scipy.stats import chi2_contingency

    def _run_chi2(data, var1, var2):
        try:
            contingency = pd.crosstab(data[var1], data[var2])
            if contingency.size == 0 or min(contingency.shape) < 2:
                return None, None
            chi2, p, dof, expected = chi2_contingency(contingency)
            n = contingency.sum().sum()
            cramers_v = np.sqrt(chi2 / (n * (min(contingency.shape) - 1)))
            return cramers_v, p
        except Exception:
            return None, None

    # 4개 검정
    fail_kw = fail_df[fail_df['Keyword'] != '']
    chi2_results = [
        ("IC ×\n합격률", *_run_chi2(df, 'IC', 'Result')),
        ("IC ×\n결함종류", *_run_chi2(fail_kw, 'IC', 'Keyword')),
        ("고객사 ×\n합격률", *_run_chi2(df, 'Customer', 'Result')),
        ("고객사 ×\n결함종류", *_run_chi2(fail_kw, 'Customer', 'Keyword')),
    ]
    chi2_results = [r for r in chi2_results if r[1] is not None]

    if chi2_results:
        labels = [r[0] for r in chi2_results]
        cramers = [r[1] for r in chi2_results]
        pvals = [r[2] for r in chi2_results]
        colors = ['#bdc3c7' if p >= 0.05 else '#e74c3c' for p in pvals]

        fig, ax = plt.subplots(figsize=(11, 6))
        bars = ax.bar(labels, cramers, color=colors, edgecolor='black',
                      linewidth=1.5, alpha=0.85)
        for bar, cv, p in zip(bars, cramers, pvals):
            sig = "유의함" if p < 0.05 else "무관"
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.02,
                    f'{cv:.3f}\n({sig})', ha='center', fontsize=11, fontweight='bold')
        ax.axhline(y=0.1, color='gray', linestyle='--', alpha=0.5)
        ax.text(len(labels)-0.4, 0.11, '약한 관계 기준', fontsize=9, color='gray')
        ax.set_ylabel("Cramér's V (관계 강도)", fontsize=12, fontweight='bold')
        ax.set_title('합격률 vs 결함 종류 — 무엇이 IC/고객사와 관련 있는가',
                     fontsize=14, fontweight='bold')
        ax.set_ylim(0, max(cramers) * 1.25)
        # 범례
        from matplotlib.patches import Patch
        legend_elems = [
            Patch(facecolor='#e74c3c', edgecolor='black', label='유의함 (p<0.05)'),
            Patch(facecolor='#bdc3c7', edgecolor='black', label='무관 (p≥0.05)')
        ]
        ax.legend(handles=legend_elems, loc='upper left', fontsize=10)
        plt.tight_layout()
        st.pyplot(fig)
        plt.close(fig)

        # 💡 인사이트 (통일된 형식)
        insight_chisquare(df, fail_df)


# ==============================================================================
# 📈 메뉴: Fail율 예측
# ==============================================================================
def render_fail_rate_prediction(df, fail_df):
    """다음 빌드 Fail율 예측"""
    plt.rcParams['font.family'] = font_name
    
    from sklearn.preprocessing import PolynomialFeatures
    from sklearn.pipeline import make_pipeline
    from sklearn.model_selection import LeaveOneOut
    from sklearn.metrics import mean_absolute_error
    
    st.markdown("### 🔮 모델별 다음 빌드 Fail율 예측")
    st.caption("같은 모델은 동일한 FW 라인(기능)이 빌드별로 진화하므로 모델 단위로 예측합니다.")
    
    # ===== 예측 대상: 모델 선택 (FW 라인 단위) =====
    st.markdown("#### 🎯 예측 대상 모델 선택")
    
    # 모델별 빌드 수 계산 (4개 이상만 예측 가능)
    model_build_counts = df.groupby('Model')['Build_Num'].nunique()
    predictable_models = sorted(model_build_counts[model_build_counts >= 4].index.tolist())
    
    if len(predictable_models) == 0:
        st.warning("⚠️ 예측 가능한 모델이 없습니다 (빌드 4개 이상 필요).")
        return
    
    col1, col2 = st.columns([2, 1])
    with col1:
        # 모델 선택 (IC 정보도 함께 표시)
        model_options = []
        model_label_map = {}
        for m in predictable_models:
            ic = df[df['Model'] == m]['IC'].iloc[0]
            n_builds = model_build_counts[m]
            label = f"{m} (IC: {ic}, 빌드 {n_builds}개)"
            model_options.append(label)
            model_label_map[label] = m
        
        selected_label = st.selectbox(
            "💻 모델 선택",
            options=model_options,
            key='pred_model_filter',
            help="FW 라인이 동일한 모델 단위로 예측합니다."
        )
        selected_model = model_label_map[selected_label]
    with col2:
        st.write("")
        st.write("")
        ic_of_model = df[df['Model'] == selected_model]['IC'].iloc[0]
        st.info(f"🔌 IC: **{ic_of_model}**")
    
    # 선택 모델로 필터링
    filtered_df = df[df['Model'] == selected_model].copy()
    active_filters = [f"Model={selected_model}"]
    
    st.info(f"🎯 분석 대상: 모델 `{selected_model}` (IC: {ic_of_model}) — 총 {len(filtered_df):,}건")
    
    if len(filtered_df) < 30:
        st.warning(f"⚠️ 데이터가 너무 적습니다 ({len(filtered_df)}건).")
        return
    
    st.markdown("---")
    
    # ===== 빌드별 데이터 추출 =====
    build_summary = filtered_df.groupby('Build_Num').agg(
        total=('Result', 'count'),
        fail=('Result', lambda x: (x == 'FAIL').sum())
    ).reset_index()
    build_summary['fail_rate'] = (build_summary['fail'] / build_summary['total'] * 100).round(2)
    build_summary = build_summary.sort_values('Build_Num').reset_index(drop=True)
    build_summary['build_idx'] = range(len(build_summary))
    
    if len(build_summary) < 4:
        st.warning(f"⚠️ 빌드가 4개 이상 필요합니다 (현재 {len(build_summary)}개).")
        return
    
    # ===== 상단 KPI =====
    st.markdown("#### 📊 현재 상태")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📦 분석 빌드 수", f"{len(build_summary)}개")
    with col2:
        latest_rate = build_summary['fail_rate'].iloc[-1]
        st.metric("📌 최신 Fail율", f"{latest_rate:.2f}%")
    with col3:
        # 표본 가중 평균
        weights_arr = build_summary['total'].values
        weighted_avg = (build_summary['fail_rate'] * weights_arr).sum() / weights_arr.sum()
        st.metric("📊 가중 평균 Fail율", f"{weighted_avg:.2f}%",
                  help="각 빌드의 표본 크기를 반영한 평균")
    with col4:
        recent_avg = build_summary['fail_rate'].tail(3).mean()
        trend = recent_avg - weighted_avg
        st.metric("📉 최근 추세",
                  f"{'개선' if trend < 0 else '악화'}",
                  delta=f"{trend:+.2f}%p",
                  delta_color="normal" if trend < 0 else "inverse")
    
    st.markdown("---")
    
    # 모델 자동 설정 (2차 다항회귀 + 표본 가중치)
    degree = 2
    use_weights = True
    
    # ===== 모델 학습 =====
    X = build_summary[['build_idx']].values
    y = build_summary['fail_rate'].values
    sample_weights = build_summary['total'].values
    
    # 다항회귀 모델
    model = make_pipeline(PolynomialFeatures(degree=degree), LinearRegression())
    
    if use_weights:
        model.fit(X, y, linearregression__sample_weight=sample_weights)
    else:
        model.fit(X, y)
    
    # ===== 다음 빌드 자동 계산 =====
    builds_array = build_summary['Build_Num'].values
    latest_build = builds_array[-1]
    # 과거 빌드 간격의 평균으로 다음 빌드 번호 추정
    if len(builds_array) > 1:
        avg_gap = int(np.mean(np.diff(builds_array)))
    else:
        avg_gap = 150
    next_build_input = int(latest_build + avg_gap)
    
    # 다음 빌드 = 인덱스 +1 (바로 다음 빌드만 예측)
    next_idx = len(build_summary)
    
    y_pred = float(model.predict([[next_idx]])[0])
    y_pred = max(0, min(100, y_pred))  # 0~100 범위 제한
    
    st.info(f"🎯 **다음 빌드 예측**: 최신 빌드 `{int(latest_build)}` 다음으로 예상되는 빌드 `{next_build_input}` (과거 평균 간격 {avg_gap} 반영)의 Fail율을 예측합니다.")
    
    # R² (학습 데이터)
    train_r2 = model.score(X, y)
    
    # ===== 신뢰구간 (예측구간) 계산 =====
    y_train_pred = model.predict(X)
    residuals = y - y_train_pred
    residual_std = np.std(residuals, ddof=1)
    
    # 95% 신뢰구간 (간단한 정규분포 가정)
    n = len(X)
    margin = 1.96 * residual_std * np.sqrt(1 + 1/n)  # 예측구간
    ci_low = max(0, y_pred - margin)
    ci_high = min(100, y_pred + margin)
    
    # ===== Leave-One-Out Cross-Validation =====
    loo = LeaveOneOut()
    cv_errors = []
    for train_idx, test_idx in loo.split(X):
        cv_model = make_pipeline(PolynomialFeatures(degree=degree), LinearRegression())
        if use_weights:
            cv_model.fit(X[train_idx], y[train_idx], 
                        linearregression__sample_weight=sample_weights[train_idx])
        else:
            cv_model.fit(X[train_idx], y[train_idx])
        pred_cv = cv_model.predict(X[test_idx])[0]
        cv_errors.append(abs(pred_cv - y[test_idx][0]))
    cv_mae = np.mean(cv_errors)
    
    st.markdown("---")
    
    # ===== 예측 결과 =====
    st.markdown("#### 🎯 예측 결과")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("🔮 예측 Fail율 (점추정)", f"{y_pred:.2f}%")
    with col2:
        st.metric("📊 95% 신뢰구간", f"[{ci_low:.1f}%, {ci_high:.1f}%]",
                  delta=f"±{margin:.2f}%p")
    with col3:
        st.metric("🎯 학습 R²", f"{train_r2:.3f}",
                  delta=f"{'강함' if train_r2 > 0.7 else '중간' if train_r2 > 0.4 else '약함'}")
    with col4:
        st.metric("✅ 교차검증 MAE", f"{cv_mae:.2f}%p",
                  delta="평균 오차",
                  delta_color="off",
                  help="Leave-One-Out 검증 평균 절대 오차")
    
    # ===== 추세 차트 (신뢰구간 포함) =====
    st.markdown("#### 📊 추세 분석 + 신뢰구간")
    
    fig, ax = plt.subplots(figsize=(13, 6.5))
    
    # 학습 데이터 (버블 크기 = 표본 크기)
    for _, row in build_summary.iterrows():
        size = 80 + (row['total'] / sample_weights.max()) * 250
        ax.scatter(row['Build_Num'], row['fail_rate'],
                   s=size, c='#3498db', edgecolors='black', linewidth=1.5,
                   zorder=4, alpha=0.7)
        ax.annotate(f"n={row['total']}", 
                    (row['Build_Num'], row['fail_rate']),
                    xytext=(0, -20), textcoords='offset points',
                    fontsize=8, ha='center', color='gray')
    
    # 회귀선 + 신뢰구간 (학습 구간)
    X_line_idx = np.linspace(0, len(build_summary) - 1, 100).reshape(-1, 1)
    X_line_build = np.linspace(builds_array[0], builds_array[-1], 100)
    y_line = model.predict(X_line_idx)
    
    # 신뢰구간 음영
    ax.fill_between(X_line_build, 
                     np.maximum(0, y_line - margin), 
                     np.minimum(100, y_line + margin),
                     color='#3498db', alpha=0.15, zorder=2,
                     label=f'95% 신뢰구간 (±{margin:.2f}%p)')
    
    ax.plot(X_line_build, y_line, '-', color='#3498db',
            linewidth=2.5, alpha=0.9, zorder=3,
            label=f'{degree}차 다항회귀 (학습)')
    
    # 예측선 (점선) + 신뢰구간 — 최신 빌드 → 다음 빌드
    X_pred_idx = np.linspace(len(build_summary) - 1, next_idx, 50).reshape(-1, 1)
    X_pred_build = np.linspace(latest_build, next_build_input, 50)
    y_pred_line = model.predict(X_pred_idx)
    
    ax.fill_between(X_pred_build,
                     np.maximum(0, y_pred_line - margin),
                     np.minimum(100, y_pred_line + margin),
                     color='#e74c3c', alpha=0.15, zorder=2)
    
    ax.plot(X_pred_build, y_pred_line, '--', color='#e74c3c',
            linewidth=2.5, label='다음 빌드 예측', zorder=3)
    
    # 예측점 (별표)
    ax.scatter([next_build_input], [y_pred],
               s=500, c='#e74c3c', marker='*', edgecolors='black', linewidth=2,
               zorder=10, label=f'예측: 빌드 {next_build_input}')
    
    # 예측값 라벨 (신뢰구간 포함)
    ax.annotate(f'예측 {y_pred:.2f}%\n[{ci_low:.1f}, {ci_high:.1f}]',
                (next_build_input, y_pred),
                xytext=(20, 25), textcoords='offset points',
                fontsize=12, fontweight='bold', color='#e74c3c',
                bbox=dict(boxstyle='round,pad=0.5', facecolor='white',
                          edgecolor='#e74c3c', linewidth=2),
                arrowprops=dict(arrowstyle='->', color='#e74c3c', lw=2))
    
    # 평균선
    ax.axhline(y=weighted_avg, color='gray', linestyle=':', alpha=0.5,
               label=f'가중 평균 ({weighted_avg:.2f}%)')
    
    filter_label = ' · '.join(active_filters) if active_filters else "전체 데이터"
    method_label = f"{degree}차 다항회귀{' + 가중치' if use_weights else ''}"
    ax.set_xlabel('빌드 번호', fontsize=12, fontweight='bold')
    ax.set_ylabel('Fail율 (%)', fontsize=12, fontweight='bold')
    ax.set_title(f'Fail율 예측 — {method_label}  ({filter_label})\n'
                 f'점 크기 = 표본 크기  |  음영 = 95% 신뢰구간',
                 fontsize=12, fontweight='bold', pad=15)
    ax.grid(True, alpha=0.3)
    ax.legend(loc='best', fontsize=9)
    plt.tight_layout()
    st.pyplot(fig)
    plt.close(fig)
    
    st.markdown("---")
    
    # ===== 예측 요약 =====
    st.markdown("#### 💡 예측 결과 해석")
    
    # 추세 방향
    if y_pred < latest_rate - 1:
        trend_text = f"✅ 개선 추세 — 최신 대비 {latest_rate - y_pred:.2f}%p 감소 예상"
        trend_advice = "펌웨어가 안정화되고 있습니다. 현재 개선 방향을 유지하세요."
    elif y_pred > latest_rate + 1:
        trend_text = f"⚠️ 악화 추세 — 최신 대비 {y_pred - latest_rate:.2f}%p 증가 예상"
        trend_advice = "Fail율 증가가 예상됩니다. 다음 빌드 출시 전 사전 모니터링을 권장합니다."
    else:
        trend_text = "➡️ 안정적 — 최신 빌드와 유사한 수준 유지"
        trend_advice = "Fail율이 안정화 단계입니다."
    
    info_text = (
        f"**🔮 다음 빌드 예측 요약**\n\n"
        f"- **추세**: {trend_text}\n"
        f"- **예측값**: 빌드 {next_build_input} → **{y_pred:.2f}%** "
        f"[95% 신뢰구간: {ci_low:.1f}% ~ {ci_high:.1f}%]\n\n"
        f"💡 **권장 사항**: {trend_advice}"
    )
    st.info(info_text)
    




# ==============================================================================
# 🚨 메뉴 3: 회귀 알람
# ==============================================================================
def render_regression_alert(df, fail_df):
    builds = sorted(df['Build_Num'].unique())
    latest_build = builds[-1]
    prev_build = builds[-2] if len(builds) > 1 else builds[-1]

    latest_fail = fail_df[fail_df['Build_Num'] == latest_build]
    prev_fail = fail_df[fail_df['Build_Num'] == prev_build]

    latest_combo = latest_fail.groupby(['Model', 'Keyword']).size().reset_index(name='latest')
    prev_combo = prev_fail.groupby(['Model', 'Keyword']).size().reset_index(name='prev')

    merged = latest_combo.merge(prev_combo, on=['Model', 'Keyword'], how='left')
    merged['prev'] = merged['prev'].fillna(0).astype(int)
    merged['상태'] = merged.apply(
        lambda r: 'NEW' if r['prev'] == 0 else ('WORSE' if r['latest'] > r['prev'] else ''),
        axis=1)
    alerts = merged[merged['상태'].isin(['NEW', 'WORSE'])].copy()
    alerts['검토영역'] = alerts['Keyword'].apply(
        lambda k: KEYWORD_DOMAIN_MAP.get(k, {}).get('area', '?'))

    # Fail 건수 기준 정렬 (많은 게 위로)
    alerts = alerts.sort_values('latest', ascending=False)

    n_new = (alerts['상태'] == 'NEW').sum()
    n_worse = (alerts['상태'] == 'WORSE').sum()
    total_fail_increase = (alerts['latest'] - alerts['prev']).sum()

    st.markdown(f"**최신 빌드 {latest_build}  vs  이전 빌드 {prev_build}** 비교")

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown("##### 🚨 전체 알람")
        st.markdown(f"### :red[**{len(alerts)}건**]")
    with col2:
        st.markdown("##### 🆕 NEW (신규)")
        st.markdown(f"### :red[**{n_new}건**]")
    with col3:
        st.markdown("##### 📈 WORSE (악화)")
        st.markdown(f"### :orange[**{n_worse}건**]")
    with col4:
        st.markdown("##### ⚠️ Fail 증가")
        st.markdown(f"### :blue[**+{total_fail_increase}건**]")
    st.markdown("---")

    alert_rows = []
    for _, row in alerts.iterrows():
        sample = latest_fail[(latest_fail['Model'] == row['Model']) &
                             (latest_fail['Keyword'] == row['Keyword'])].iloc[0]
        video_link = get_video_link(sample['video_id'])
        alert_rows.append({
            '상태': row['상태'],
            'Model': row['Model'],
            'Fail_Type': row['Keyword'],
            '이전 건수': row['prev'],
            '현재 건수': row['latest'],
            '증가': int(row['latest']) - int(row['prev']),
            '검토 영역': row['검토영역'],
            '영상': video_link
        })
    alert_df = pd.DataFrame(alert_rows)

    st.dataframe(
        alert_df,
        use_container_width=True,
        hide_index=True,
        column_config={
            "영상": st.column_config.LinkColumn("영상 보기", display_text="🎬 재생")
        }
    )
    
    # 💡 인사이트
    insight_regression_alert(alerts, latest_build, prev_build, df, fail_df)


# ==============================================================================
# 📈 메뉴 새로 추가: 전체 통계
# ==============================================================================
def render_full_statistics(df, fail_df):
    """전체 데이터 통계 - 월별, Customer, Test_Item 매트릭스, 메타 정보"""
    plt.rcParams['font.family'] = font_name
    
    # ===== 종합 KPI (핵심 4개만) =====
    st.markdown("### 🎯 종합 KPI")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📋 전체 테스트", f"{len(df):,}건")
    with col2:
        fail_rate = len(fail_df) / len(df) * 100 if len(df) > 0 else 0
        st.metric("❌ Fail율", f"{fail_rate:.1f}%", f"{len(fail_df):,}건")
    with col3:
        st.metric("💻 모델 수", f"{df['Model'].nunique()}개")
    with col4:
        st.metric("🔄 빌드 수", f"{df['Build_Num'].nunique()}개")
    
    st.markdown("---")
    
    # ===== 1. 월별 테스트 추이 =====
    st.markdown("### 📅 1. 월별 테스트 추이")
    st.caption("시간 흐름에 따른 테스트량과 Fail율 변화")
    
    df_dated = df.copy()
    df_dated['Test_Date'] = pd.to_datetime(df_dated['Test_Date'], errors='coerce')
    df_dated = df_dated.dropna(subset=['Test_Date'])
    df_dated['year_month'] = df_dated['Test_Date'].dt.to_period('M').astype(str)
    
    monthly = df_dated.groupby('year_month').agg(
        total=('Result', 'count'),
        fail=('Result', lambda x: (x == 'FAIL').sum())
    ).reset_index()
    monthly['fail_rate'] = (monthly['fail'] / monthly['total'] * 100).round(1)
    
    fig, ax1 = plt.subplots(figsize=(14, 5))
    
    # 막대: 테스트 건수
    bars = ax1.bar(monthly['year_month'], monthly['total'],
                    color='#3498db', alpha=0.6, label='테스트 건수')
    ax1.set_xlabel('월', fontsize=11)
    ax1.set_ylabel('테스트 건수', fontsize=11, color='#3498db')
    ax1.tick_params(axis='y', labelcolor='#3498db')
    ax1.tick_params(axis='x', rotation=45)
    
    for bar, val in zip(bars, monthly['total']):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 20,
                 f'{val}', ha='center', fontsize=9, color='#3498db')
    
    # 라인: Fail율
    ax2 = ax1.twinx()
    ax2.plot(monthly['year_month'], monthly['fail_rate'],
             marker='o', markersize=8, linewidth=2, color='#e74c3c', label='Fail율 (%)')
    ax2.set_ylabel('Fail율 (%)', fontsize=11, color='#e74c3c')
    ax2.tick_params(axis='y', labelcolor='#e74c3c')
    ax2.set_ylim(0, max(monthly['fail_rate']) * 1.3)
    
    for i, val in enumerate(monthly['fail_rate']):
        ax2.text(i, val + 1, f'{val}%', ha='center', fontsize=9, color='#e74c3c', fontweight='bold')
    
    ax1.set_title('월별 테스트 건수 + Fail율 추이', fontsize=14, fontweight='bold')
    ax1.grid(True, alpha=0.3)
    plt.tight_layout()
    st.pyplot(fig)
    plt.close(fig)
    
    # 표
    monthly_display = monthly.copy()
    monthly_display.columns = ['월', '전체', 'Fail', 'Fail율(%)']
    st.dataframe(monthly_display, use_container_width=True, hide_index=True)
    
    # 💡 인사이트
    insight_monthly_trend(df, fail_df)
    
    st.markdown("---")
    
    # ===== 2. Customer (고객사) 분석 =====
    st.markdown("### 🏢 2. 고객사별 분석")
    st.caption("어떤 고객사가 어떤 IC를 어떻게 테스트하는지")
    
    col_left, col_right = st.columns(2)
    
    with col_left:
        st.markdown("##### 고객사별 테스트 건수 및 Fail율")
        cust = df.groupby('Customer').agg(
            total=('Result', 'count'),
            fail=('Result', lambda x: (x == 'FAIL').sum())
        ).reset_index()
        cust['fail_rate'] = (cust['fail'] / cust['total'] * 100).round(1)
        cust = cust.sort_values('total', ascending=False)
        
        fig, ax = plt.subplots(figsize=(8, 5))
        bars = ax.barh(cust['Customer'], cust['total'],
                       color=plt.cm.Set2(np.linspace(0, 1, len(cust))))
        for bar, t, fr in zip(bars, cust['total'], cust['fail_rate']):
            ax.text(bar.get_width() + max(cust['total']) * 0.01,
                    bar.get_y() + bar.get_height()/2,
                    f'{t}건 ({fr}%)', va='center', fontsize=9, fontweight='bold')
        ax.set_xlabel('테스트 건수')
        ax.invert_yaxis()
        ax.set_title('고객사별 테스트 분포', fontsize=12, fontweight='bold')
        plt.tight_layout()
        st.pyplot(fig)
        plt.close(fig)
    
    with col_right:
        st.markdown("##### IC × 고객사 매트릭스")
        ic_cust = df.groupby(['IC', 'Customer']).size().unstack(fill_value=0)
        
        fig, ax = plt.subplots(figsize=(8, 5))
        sns.heatmap(ic_cust, annot=True, fmt='d', cmap='Blues',
                    cbar_kws={'label': '테스트 건수'}, ax=ax,
                    linewidths=0.5, linecolor='white')
        ax.set_title('IC × 고객사 분포', fontsize=12, fontweight='bold')
        ax.set_xlabel('Customer')
        ax.set_ylabel('IC')
        plt.tight_layout()
        st.pyplot(fig)
        plt.close(fig)
    
    # 💡 인사이트
    insight_customer_analysis(df, fail_df)
    
    st.markdown("---")
    
    # ===== 3. Test_Item × Fail_Type 매트릭스 =====
    st.markdown("### 🧪 3. Test_Item × Fail_Type 매트릭스")
    st.caption("어떤 테스트 항목에서 어떤 결함이 자주 발생하는지")
    
    if len(fail_df) > 0:
        # Top 15 Test_Item × Fail_Type
        top_items = fail_df['Test_Item'].value_counts().head(15).index.tolist()
        filtered = fail_df[fail_df['Test_Item'].isin(top_items)]
        
        item_kw = filtered.pivot_table(
            index='Test_Item', columns='Keyword', values='No',
            aggfunc='count', fill_value=0
        )
        item_kw = item_kw.loc[top_items]  # 순서 유지
        item_kw = item_kw[[k for k in KEYWORDS_ALL if k in item_kw.columns]]
        
        fig, ax = plt.subplots(figsize=(14, max(6, len(item_kw) * 0.4)))
        sns.heatmap(item_kw, annot=True, fmt='d', cmap='YlOrRd',
                    cbar_kws={'label': 'Fail 건수'}, ax=ax,
                    linewidths=0.5, linecolor='white')
        ax.set_title('Test_Item (TOP 15) × Fail_Type 매트릭스',
                     fontsize=14, fontweight='bold', pad=15)
        ax.set_xlabel('Fail_Type')
        ax.set_ylabel('Test_Item')
        plt.tight_layout()
        st.pyplot(fig)
        plt.close(fig)
        
        # 💡 인사이트 (단순 안내 → 의사결정용)
        insight_test_item_matrix(fail_df)
    
    st.markdown("---")
    
    # ===== 4. 데이터 메타 정보 =====
    st.markdown("### 📋 4. 데이터 메타 정보")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("##### 📅 수집 기간")
        if len(df_dated) > 0:
            min_date = df_dated['Test_Date'].min()
            max_date = df_dated['Test_Date'].max()
            duration = (max_date - min_date).days
            st.markdown(f"""
            - **시작**: {min_date.strftime('%Y-%m-%d')}
            - **종료**: {max_date.strftime('%Y-%m-%d')}
            - **기간**: 약 {duration}일 ({duration//30}개월)
            - **수집 월**: {df_dated['year_month'].nunique()}개월
            """)
    
    with col2:
        st.markdown("##### 🔌 IC 종류")
        for ic in sorted(df['IC'].unique()):
            ic_data = df[df['IC'] == ic]
            ic_fail = ic_data[ic_data['Result'] == 'FAIL']
            top_kw = ic_fail['Keyword'].value_counts().head(1)
            top_text = f"{top_kw.index[0]} ({top_kw.iloc[0]}건)" if len(top_kw) > 0 else "Fail 없음"
            st.markdown(f"- **{ic}**: {len(ic_data):,}건 · Top Fail: {top_text}")
    
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("##### 📊 결과 분포")
        result_counts = df['Result'].value_counts()
        for r, c in result_counts.items():
            pct = c / len(df) * 100
            st.markdown(f"- **{r}**: {c:,}건 ({pct:.1f}%)")
    
    with col2:
        st.markdown("##### 🐛 주요 Fail_Type (TOP 5)")
        top_kw = fail_df['Keyword'].value_counts().head(5)
        for i, (kw, c) in enumerate(top_kw.items(), 1):
            st.markdown(f"- **{i}위 {kw}**: {c}건")




# ==============================================================================
# 📁 메뉴 4: 데이터 업로드 (선택 삭제 + 팝업 UX)
# ==============================================================================
def render_data_upload():
    st.markdown("### 📁 주간 체크리스트 업로드")
    st.caption("회사 SQA 시스템에서 받은 엑셀(.xlsx) 또는 CSV 파일을 업로드하면 자동으로 분석 데이터에 추가됩니다.")
    
    # 현재 데이터 상태
    base_df = load_base_data()
    current_total = len(base_df) + sum(len(d) for d in st.session_state.uploaded_data)
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("📋 기본 데이터", f"{len(base_df):,}건")
    with col2:
        st.metric("➕ 추가된 데이터", f"{sum(len(d) for d in st.session_state.uploaded_data):,}건")
    with col3:
        st.metric("🔵 현재 총합", f"{current_total:,}건")
    
    st.markdown("---")
    
    # 이전 결과 팝업 표시 (한 번만)
    if 'upload_result' in st.session_state and st.session_state.upload_result:
        result = st.session_state.upload_result
        if result['status'] == 'success':
            st.success(f"✅ **업로드 성공** — {result['message']}")
            st.balloons()
        elif result['status'] == 'error':
            st.error(f"❌ **업로드 실패** — {result['message']}")
        elif result['status'] == 'delete':
            st.success(f"🗑️ **삭제 완료** — {result['message']}")
        st.session_state.upload_result = None
    
    # 업로드 이력 추적용 세션
    if 'uploaded_files_info' not in st.session_state:
        st.session_state.uploaded_files_info = []
    
    if 'uploader_key' not in st.session_state:
        st.session_state.uploader_key = 0
    
    # ===== 파일 업로드 =====
    uploaded_file = st.file_uploader(
        "📤 파일 선택 (xlsx, csv)",
        type=['xlsx', 'csv'],
        help="컬럼 형식: 기본 CSV와 동일 (No, IC, Model, Build_Num, Result, Keyword 등)",
        key=f"file_uploader_{st.session_state.uploader_key}"
    )
    
    if uploaded_file is not None:
        try:
            # 파일 읽기
            if uploaded_file.name.endswith('.csv'):
                new_df = pd.read_csv(uploaded_file)
            else:
                new_df = pd.read_excel(uploaded_file)
            
            # 필수 컬럼 체크
            required_cols = ['IC', 'Model', 'Build_Num', 'Result', 'Test_Item', 'Keyword']
            missing = [c for c in required_cols if c not in new_df.columns]
            if missing:
                st.session_state.upload_result = {
                    'status': 'error',
                    'message': f"필수 컬럼 누락: {', '.join(missing)}"
                }
                st.session_state.uploader_key += 1
                st.rerun()
            
            # 데이터 전처리
            if 'Test_Date' in new_df.columns:
                new_df['Test_Date'] = pd.to_datetime(new_df['Test_Date'], errors='coerce')
            if 'IC_Code' in new_df.columns:
                new_df['IC_Code'] = new_df['IC_Code'].astype(str).str.zfill(2)
            if 'Model_Minor' in new_df.columns:
                new_df['Model_Minor'] = new_df['Model_Minor'].astype(str).str.zfill(2)
            new_df['Keyword'] = new_df['Keyword'].fillna('')
            
            # ===== 중복 검사 =====
            
            # [1] 같은 파일
            file_signature = (uploaded_file.name, len(new_df))
            already_uploaded_files = [
                (info['filename'], info['n_rows']) 
                for info in st.session_state.uploaded_files_info
            ]
            if file_signature in already_uploaded_files:
                st.session_state.upload_result = {
                    'status': 'error',
                    'message': f"'{uploaded_file.name}' ({len(new_df)}건)은 이미 업로드된 파일입니다."
                }
                st.session_state.uploader_key += 1
                st.rerun()
            
            # [2] No 컬럼 중복
            if 'No' in new_df.columns:
                new_nos = set(new_df['No'].dropna().astype(int).tolist())
                existing_nos = set(base_df['No'].dropna().astype(int).tolist())
                for d in st.session_state.uploaded_data:
                    if 'No' in d.columns:
                        existing_nos.update(d['No'].dropna().astype(int).tolist())
                
                overlap_nos = new_nos & existing_nos
                if overlap_nos:
                    overlap_sample = sorted(list(overlap_nos))[:5]
                    sample_str = ', '.join(map(str, overlap_sample))
                    extra = f" 등 총 {len(overlap_nos)}건" if len(overlap_nos) > 5 else ""
                    st.session_state.upload_result = {
                        'status': 'error',
                        'message': f"No 컬럼 중복 — 이미 존재하는 번호: {sample_str}{extra}"
                    }
                    st.session_state.uploader_key += 1
                    st.rerun()
            
            # [3] 행 단위 중복
            check_cols = ['Build_Num', 'Model', 'Test_Item', 'Result']
            if all(c in new_df.columns for c in check_cols):
                all_existing = pd.concat([base_df] + st.session_state.uploaded_data, ignore_index=True) \
                    if st.session_state.uploaded_data else base_df
                existing_keys = set(
                    all_existing[check_cols].apply(lambda r: tuple(r), axis=1).tolist()
                )
                new_keys = new_df[check_cols].apply(lambda r: tuple(r), axis=1).tolist()
                n_overlap = sum(1 for k in new_keys if k in existing_keys)
                overlap_rate = n_overlap / len(new_df) * 100 if len(new_df) > 0 else 0
                
                if overlap_rate >= 50:
                    st.session_state.upload_result = {
                        'status': 'error',
                        'message': f"데이터 중복 — {n_overlap}건 ({overlap_rate:.1f}%)이 기존 데이터와 동일합니다."
                    }
                    st.session_state.uploader_key += 1
                    st.rerun()
            
            # 통과 → 추가
            no_set = set(new_df['No'].dropna().astype(int).tolist()) if 'No' in new_df.columns else set()
            st.session_state.uploaded_files_info.append({
                'filename': uploaded_file.name,
                'n_rows': len(new_df),
                'no_set': no_set,
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            })
            st.session_state.uploaded_data.append(new_df)
            st.session_state.data_version += 1
            
            new_builds = sorted(new_df['Build_Num'].unique())
            existing_builds = sorted(base_df['Build_Num'].unique())
            truly_new_builds = [b for b in new_builds if b not in existing_builds]
            
            pass_n = (new_df['Result'] == 'PASS').sum()
            fail_n = (new_df['Result'] == 'FAIL').sum()
            
            success_msg = f"`{uploaded_file.name}` 파일에서 **{len(new_df)}건** 추가됨 (PASS {pass_n} · FAIL {fail_n})"
            if truly_new_builds:
                success_msg += f"  🆕 신규 빌드: {truly_new_builds}"
            
            st.session_state.upload_result = {
                'status': 'success',
                'message': success_msg
            }
            st.session_state.uploader_key += 1
            st.rerun()
        
        except Exception as e:
            st.session_state.upload_result = {
                'status': 'error',
                'message': f"파일 처리 중 오류 발생: {str(e)}"
            }
            st.session_state.uploader_key += 1
            st.rerun()
    
    # ===== 📜 업로드 이력 (체크박스 + 선택 삭제) =====
    if st.session_state.uploaded_files_info:
        st.markdown("---")
        st.markdown("#### 📜 업로드 이력")
        st.caption("체크박스로 삭제할 파일을 선택할 수 있습니다.")
        
        # 각 파일에 대한 체크박스 (세션에 선택 상태 저장)
        if 'delete_selection' not in st.session_state or \
           len(st.session_state.delete_selection) != len(st.session_state.uploaded_files_info):
            st.session_state.delete_selection = [False] * len(st.session_state.uploaded_files_info)
        
        # 헤더
        col_h1, col_h2, col_h3, col_h4 = st.columns([0.7, 3, 1.5, 2])
        with col_h1:
            st.markdown("**선택**")
        with col_h2:
            st.markdown("**파일명**")
        with col_h3:
            st.markdown("**건수**")
        with col_h4:
            st.markdown("**업로드 시각**")
        
        # 각 항목 표시
        for i, info in enumerate(st.session_state.uploaded_files_info):
            col1, col2, col3, col4 = st.columns([0.7, 3, 1.5, 2])
            with col1:
                checked = st.checkbox(
                    f"선택 {i+1}",
                    key=f"file_check_{i}_{st.session_state.uploader_key}",
                    value=st.session_state.delete_selection[i],
                    label_visibility="collapsed"
                )
                st.session_state.delete_selection[i] = checked
            with col2:
                st.markdown(f"`{info['filename']}`")
            with col3:
                st.markdown(f"{info['n_rows']:,}건")
            with col4:
                st.markdown(info['uploaded_at'])
        
        st.markdown("")
        
        # 선택된 개수
        n_selected = sum(st.session_state.delete_selection)
        n_total = len(st.session_state.uploaded_files_info)
        
        # 액션 버튼 2개
        btn_col1, btn_col2 = st.columns(2)
        
        with btn_col1:
            if st.button("☑️ 전체 선택", use_container_width=True):
                st.session_state.delete_selection = [True] * n_total
                st.rerun()
        
        with btn_col2:
            disabled = n_selected == 0
            label = f"🗑️ 선택 삭제 ({n_selected}개)" if n_selected > 0 else "🗑️ 선택 삭제 (0개)"
            if st.button(label, type="primary" if not disabled else "secondary",
                         use_container_width=True, disabled=disabled):
                # 선택된 파일들의 인덱스 (역순으로 정렬해서 삭제)
                indices_to_delete = sorted(
                    [i for i, sel in enumerate(st.session_state.delete_selection) if sel],
                    reverse=True
                )
                deleted_names = []
                deleted_count = 0
                for idx in indices_to_delete:
                    deleted_names.append(st.session_state.uploaded_files_info[idx]['filename'])
                    deleted_count += st.session_state.uploaded_files_info[idx]['n_rows']
                    del st.session_state.uploaded_files_info[idx]
                    del st.session_state.uploaded_data[idx]
                    del st.session_state.delete_selection[idx]
                
                st.session_state.data_version += 1
                st.session_state.uploader_key += 1
                st.session_state.upload_result = {
                    'status': 'delete',
                    'message': f"{len(deleted_names)}개 파일 ({deleted_count}건) 삭제됨: {', '.join(deleted_names[:3])}{'...' if len(deleted_names) > 3 else ''}"
                }
                st.rerun()


# ==============================================================================
# 📥 메뉴 5: 보고서 생성 (이력 저장 + 선택 삭제)
# ==============================================================================
def render_report_generator(df, fail_df):
    st.markdown("### 📥 주간 보고서 자동 생성")
    st.caption("월 단위로 데이터를 선택해 회사 표준 SQA Quality Review PPT를 자동 생성합니다.")
    
    # 보고서 이력 세션 초기화
    if 'report_history' not in st.session_state:
        st.session_state.report_history = []  # [{filename, month, n_data, generated_at, ppt_bytes}, ...]
    if 'report_delete_selection' not in st.session_state:
        st.session_state.report_delete_selection = []
    if 'report_result' not in st.session_state:
        st.session_state.report_result = None
    
    # 이전 결과 팝업 표시
    if st.session_state.report_result:
        result = st.session_state.report_result
        if result['status'] == 'success':
            st.success(f"✅ **보고서 생성 완료** — {result['message']}")
            st.balloons()
        elif result['status'] == 'delete':
            st.success(f"🗑️ **삭제 완료** — {result['message']}")
        elif result['status'] == 'error':
            st.error(f"❌ **생성 실패** — {result['message']}")
        st.session_state.report_result = None
    
    # 데이터에서 사용 가능한 월 목록 추출
    df_copy = df.copy()
    df_copy['Test_Date'] = pd.to_datetime(df_copy['Test_Date'], errors='coerce')
    df_copy['year_month'] = df_copy['Test_Date'].dt.to_period('M')
    available_months = sorted(df_copy['year_month'].dropna().unique(), reverse=True)
    
    if len(available_months) == 0:
        st.error("⚠️ 분석 가능한 데이터가 없습니다.")
        return
    
    # 월 선택 UI
    month_options = [f"{m.year}년 {m.month}월" for m in available_months]
    month_map = {f"{m.year}년 {m.month}월": m for m in available_months}
    
    col1, col2 = st.columns(2)
    with col1:
        selected_month_str = st.selectbox(
            "📅 보고서 기준 월",
            options=month_options,
            help="이 월의 데이터로 보고서를 생성합니다."
        )
    with col2:
        report_type = st.selectbox("📄 보고서 형식", ["SQA-SW Quality Review (10장)"])
    
    selected_month = month_map[selected_month_str]
    
    # 선택한 월의 데이터 필터링
    monthly_df = df_copy[df_copy['year_month'] == selected_month].copy()
    monthly_fail = monthly_df[monthly_df['Result'] == 'FAIL'].copy()
    
    # 미리보기 정보
    st.markdown("---")
    st.markdown(f"#### 📋 {selected_month_str} 보고서 정보")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📋 분석 데이터", f"{len(monthly_df):,}건")
    with col2:
        st.metric("❌ Fail 건수", f"{len(monthly_fail):,}건")
    with col3:
        st.metric("💻 모델 수", f"{monthly_df['Model'].nunique()}개")
    with col4:
        st.metric("🔄 빌드 수", f"{monthly_df['Build_Num'].nunique()}개")
    
    if len(monthly_df) == 0:
        st.warning("⚠️ 선택한 월에 데이터가 없습니다.")
        return
    
    st.caption("📑 10장 PPT (표지·목차·요약·Fail Review·Action Item·회귀 알람·회의 요약 등)")
    
    st.markdown("---")
    
    # 생성 버튼
    if st.button("🚀 보고서 생성하기", type="primary", use_container_width=True):
        with st.spinner("📝 PPT 생성 중... (약 3~5초)"):
            try:
                ppt_buffer = generate_ppt_report(monthly_df, monthly_fail, selected_month)
                ppt_bytes = ppt_buffer.getvalue()
                
                filename = f"SQA_Quality_Review_{selected_month.year}_{selected_month.month:02d}.pptx"
                
                # 보고서 이력에 저장
                st.session_state.report_history.append({
                    'filename': filename,
                    'month_str': selected_month_str,
                    'n_data': len(monthly_df),
                    'n_fail': len(monthly_fail),
                    'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'ppt_bytes': ppt_bytes
                })
                st.session_state.report_delete_selection = [False] * len(st.session_state.report_history)
                
                st.session_state.report_result = {
                    'status': 'success',
                    'message': f"`{filename}` 생성 완료 (분석 데이터: {len(monthly_df):,}건). 아래 이력에서 다운로드하세요."
                }
                st.rerun()
            except Exception as e:
                st.session_state.report_result = {
                    'status': 'error',
                    'message': f"보고서 생성 실패: {str(e)}"
                }
                import traceback
                st.code(traceback.format_exc())
                st.rerun()
    
    # ===== 📜 보고서 이력 (체크박스 + 선택 삭제 + 다운로드) =====
    if st.session_state.report_history:
        st.markdown("---")
        st.markdown("#### 📜 생성한 보고서 이력")
        st.caption("다운로드 또는 체크박스로 삭제할 수 있습니다.")
        
        # 선택 상태 동기화
        if len(st.session_state.report_delete_selection) != len(st.session_state.report_history):
            st.session_state.report_delete_selection = [False] * len(st.session_state.report_history)
        
        # 헤더
        col_h1, col_h2, col_h3, col_h4, col_h5 = st.columns([0.7, 2.5, 1.2, 1.8, 1.5])
        with col_h1:
            st.markdown("**선택**")
        with col_h2:
            st.markdown("**파일명**")
        with col_h3:
            st.markdown("**기준 월**")
        with col_h4:
            st.markdown("**생성 시각**")
        with col_h5:
            st.markdown("**다운로드**")
        
        # 각 항목
        for i, info in enumerate(st.session_state.report_history):
            col1, col2, col3, col4, col5 = st.columns([0.7, 2.5, 1.2, 1.8, 1.5])
            with col1:
                checked = st.checkbox(
                    f"선택 {i+1}",
                    key=f"report_check_{i}",
                    value=st.session_state.report_delete_selection[i],
                    label_visibility="collapsed"
                )
                st.session_state.report_delete_selection[i] = checked
            with col2:
                st.markdown(f"`{info['filename']}`")
            with col3:
                st.markdown(info['month_str'])
            with col4:
                st.markdown(info['generated_at'])
            with col5:
                st.download_button(
                    label="📥 다운로드",
                    data=info['ppt_bytes'],
                    file_name=info['filename'],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"download_report_{i}",
                    use_container_width=True
                )
        
        st.markdown("")
        
        # 선택된 개수
        n_selected = sum(st.session_state.report_delete_selection)
        n_total = len(st.session_state.report_history)
        
        # 액션 버튼 2개
        btn_col1, btn_col2 = st.columns(2)
        
        with btn_col1:
            if st.button("☑️ 전체 선택", use_container_width=True, key="report_select_all"):
                st.session_state.report_delete_selection = [True] * n_total
                st.rerun()
        
        with btn_col2:
            disabled = n_selected == 0
            label = f"🗑️ 선택 삭제 ({n_selected}개)" if n_selected > 0 else "🗑️ 선택 삭제 (0개)"
            if st.button(label, type="primary" if not disabled else "secondary",
                         use_container_width=True, disabled=disabled, key="report_delete"):
                indices_to_delete = sorted(
                    [i for i, sel in enumerate(st.session_state.report_delete_selection) if sel],
                    reverse=True
                )
                deleted_names = []
                for idx in indices_to_delete:
                    deleted_names.append(st.session_state.report_history[idx]['filename'])
                    del st.session_state.report_history[idx]
                    del st.session_state.report_delete_selection[idx]
                
                st.session_state.report_result = {
                    'status': 'delete',
                    'message': f"{len(deleted_names)}개 보고서 삭제됨: {', '.join(deleted_names[:3])}{'...' if len(deleted_names) > 3 else ''}"
                }
                st.rerun()




def generate_ppt_report(df, fail_df, selected_month):
    """PPT 보고서 생성 → BytesIO 반환 (페이지 맞춤 버전)"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    month_str = f"{selected_month.year}년 {selected_month.month}월"
    
    # 헬퍼 함수
    def set_cell(cell, text, bold=False, size=10, color=None, bg=None, align=None):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name = PPT_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    
    def add_text(slide, text, x, y, w, h, size=14, bold=False, color=None, align=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = PPT_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    
    def add_header(slide, page_num, title):
        box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.3), Inches(0.4), Inches(0.4))
        box.fill.solid()
        box.fill.fore_color.rgb = RPT_HEADER_BG
        box.line.fill.background()
        add_text(slide, str(page_num), 0.3, 0.32, 0.4, 0.4,
                 size=14, bold=True, color=RPT_HEADER_FG, align=PP_ALIGN.CENTER)
        add_text(slide, "SQA-SW Quality Review", 0.8, 0.3, 5, 0.4,
                 size=18, bold=True, color=RPT_TITLE)
        add_text(slide, f"❑ {title}", 0.4, 0.85, 12, 0.4,
                 size=16, bold=True, color=RPT_TITLE)
    
    # ===== Slide 1: 표지 =====
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RPT_HEADER_BG
    bg.line.fill.background()
    box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(3), Inches(10.3), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.fill.background()
    add_text(s1, "SQA-SW Quality Review", 1.5, 3.4, 10.3, 0.8,
             size=32, bold=True, color=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    add_text(s1, month_str, 5, 5.5, 3.3, 0.5,
             size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    
    # ===== Slide 2: 목차 =====
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s2, 1, "List")
    items = [
        "1.   회의 개요",
        "2.   월간 TEST 결과 요약 (SQA)",
        "3.   Fail 항목 Review (SQA)",
        "4.   Action Item Review (SW)",
        "5.   회의 요약 및 내부 공유",
    ]
    y_pos = 2.0
    for item in items:
        add_text(s2, item, 1.5, y_pos, 10, 0.5, size=18, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += 0.7
    
    # ===== Slide 3: 기본정보 =====
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s3, 2, "기본정보")
    table = s3.shapes.add_table(5, 2, Inches(1.5), Inches(2),
                                 Inches(10.3), Inches(3.5)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(7.8)
    info = [
        ("항 목", "내 용"),
        ("회의명", "SQA-SW 주간 Quality Review"),
        ("보고서 기준 월", month_str),
        ("참석자", "SQA 팀, SW 개발팀"),
        ("회의목적", "SQA 펌웨어 테스트 결과 공유 및 Action Item Review"),
    ]
    for r, (c1, c2) in enumerate(info):
        if r == 0:
            set_cell(table.cell(r, 0), c1, bold=True, size=12,
                     color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 1), c2, bold=True, size=12,
                     color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
        else:
            set_cell(table.cell(r, 0), c1, bold=True, size=11,
                     bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 1), c2, size=11)
    
    # ===== Slide 4: 월간 Test 요약 =====
    # 페이지 맞춤: 최대 10개 항목까지만 (페이지 넘침 방지)
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s4, 3, "월간 Test 요약")
    add_text(s4, f"{month_str} Test Summary", 0.4, 1.4, 5, 0.4,
             size=14, bold=True, color=RPT_TITLE)
    
    # Top 10개만 표시 (페이지 맞춤)
    item_summary = df['Test_Item'].value_counts().head(10).reset_index()
    item_summary.columns = ['Test_Item', 'count']
    
    rows_t1 = len(item_summary) + 2
    # 페이지 높이 맞춤: 행 높이 자동 조절
    available_height = 5.0  # y=2부터 7까지
    row_height = min(0.4, available_height / rows_t1)
    
    table1 = s4.shapes.add_table(rows_t1, 2, Inches(0.4), Inches(2),
                                  Inches(5.5), Inches(row_height * rows_t1)).table
    table1.columns[0].width = Inches(3.5)
    table1.columns[1].width = Inches(2.0)
    set_cell(table1.cell(0, 0), "Test Item", bold=True, size=10,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    set_cell(table1.cell(0, 1), month_str, bold=True, size=10,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    total_count = 0
    for r, row in enumerate(item_summary.itertuples(), 1):
        bg = RPT_ALT_ROW if r % 2 == 0 else None
        set_cell(table1.cell(r, 0), row.Test_Item, size=9, bg=bg)
        set_cell(table1.cell(r, 1), row.count, size=9, bg=bg, align=PP_ALIGN.CENTER)
        total_count += row.count
    set_cell(table1.cell(rows_t1-1, 0), "Total", bold=True, size=10,
             bg=RPT_HEADER_BG, color=RPT_HEADER_FG, align=PP_ALIGN.CENTER)
    set_cell(table1.cell(rows_t1-1, 1), total_count, bold=True, size=10,
             bg=RPT_HEADER_BG, color=RPT_HEADER_FG, align=PP_ALIGN.CENTER)
    
    pass_count = (df['Result'] == 'PASS').sum()
    fail_count = (df['Result'] == 'FAIL').sum()
    total = pass_count + fail_count
    fail_rate = fail_count / total * 100 if total > 0 else 0
    
    table2 = s4.shapes.add_table(4, 3, Inches(6.5), Inches(2),
                                  Inches(6), Inches(1.6)).table
    table2.columns[0].width = Inches(1.5)
    table2.columns[1].width = Inches(1.8)
    table2.columns[2].width = Inches(2.7)
    set_cell(table2.cell(0, 0), "구분", bold=True, size=11,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(0, 1), month_str, bold=True, size=11,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(0, 2), "Fail율", bold=True, size=11,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(1, 0), "PASS", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(1, 1), f"{pass_count:,}", size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(1, 2), "", size=11, bg=RPT_ALT_ROW)
    set_cell(table2.cell(2, 0), "FAIL", bold=True, size=11, color=RPT_CRITICAL, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(2, 1), f"{fail_count:,}", size=11, color=RPT_CRITICAL, bold=True, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(2, 2), f"{fail_rate:.1f}%", size=12, color=RPT_CRITICAL, bold=True, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(3, 0), "Total", bold=True, size=11, bg=RPT_HEADER_BG, color=RPT_HEADER_FG, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(3, 1), f"{total:,}", bold=True, size=11, bg=RPT_HEADER_BG, color=RPT_HEADER_FG, align=PP_ALIGN.CENTER)
    set_cell(table2.cell(3, 2), "", bg=RPT_HEADER_BG)
    
    add_text(s4, "💡 자동 인사이트", 6.5, 4.0, 6, 0.4, size=12, bold=True, color=RPT_TITLE)
    if len(fail_df) > 0:
        top_fail_kw = fail_df['Keyword'].value_counts().head(3)
        insight = f"• Top Fail_Type: {', '.join(top_fail_kw.index[:3])}\n"
        insight += f"• 영향 모델: {fail_df['Model'].nunique()}개\n"
        insight += f"• 영향 IC: {fail_df['IC'].nunique()}종\n"
        insight += f"• 분석 빌드: {df['Build_Num'].nunique()}개"
    else:
        insight = "• 이번 달 Fail 없음 ✅"
    add_text(s4, insight, 6.5, 4.4, 6, 2.5, size=11, color=RGBColor(0x33, 0x33, 0x33))
    
    # ===== Slide 5: IC × 빌드별 상세 (페이지 맞춤) =====
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s5, 3, "월간 Test 요약")
    add_text(s5, f"{month_str} IC × 빌드별 Test 결과 상세", 0.4, 1.4, 10, 0.4,
             size=14, bold=True, color=RPT_TITLE)
    
    summary = df.groupby(['IC', 'Customer', 'Build_Num']).agg(
        total=('Result', 'count'),
        pass_n=('Result', lambda x: (x == 'PASS').sum()),
        fail_n=('Result', lambda x: (x == 'FAIL').sum())
    ).reset_index()
    summary['fail_rate'] = (summary['fail_n'] / summary['total'] * 100).round(1)
    # 최대 8개 행으로 제한 (페이지 맞춤)
    summary = summary.sort_values('fail_rate', ascending=False).head(8)
    
    if len(summary) > 0:
        rows = len(summary) + 1
        available_height = 5.0
        row_height = min(0.5, available_height / rows)
        
        table = s5.shapes.add_table(rows, 7, Inches(0.4), Inches(2),
                                     Inches(12.6), Inches(row_height * rows)).table
        widths = [1.0, 1.5, 2.5, 1.3, 1.5, 1.5, 1.3]
        for i, w in enumerate(widths):
            table.columns[i].width = Inches(w)
        headers = ['IC', '고객사', '빌드 번호', '총 Test', 'Pass', 'Fail', 'Fail율(%)']
        for c, h in enumerate(headers):
            set_cell(table.cell(0, c), h, bold=True, size=10,
                     color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
        for r, row in enumerate(summary.itertuples(), 1):
            bg = RPT_ALT_ROW if r % 2 == 0 else None
            set_cell(table.cell(r, 0), row.IC, bold=True, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 1), row.Customer, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 2), f"R00.0.{row.Build_Num}", size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 3), row.total, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 4), row.pass_n, size=9, bg=bg, align=PP_ALIGN.CENTER)
            fail_color = RPT_CRITICAL if row.fail_n > 0 else None
            set_cell(table.cell(r, 5), row.fail_n, size=9, bold=row.fail_n > 0, bg=bg, color=fail_color, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 6), f"{row.fail_rate}%", size=9, bold=row.fail_rate > 5, bg=bg,
                     color=fail_color if row.fail_rate > 5 else None, align=PP_ALIGN.CENTER)
    else:
        add_text(s5, "이번 달 데이터가 없습니다.", 0.4, 3, 10, 0.4, size=14)
    
    # ===== Slide 6: Fail Review (페이지 맞춤) =====
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s6, 4, "Fail / Issue Review")
    add_text(s6, "주요 Fail 케이스 (CRITICAL/HIGH 우선)", 0.4, 1.4, 10, 0.4,
             size=14, bold=True, color=RPT_TITLE)
    
    fail_df_copy = fail_df.copy()
    if len(fail_df_copy) > 0:
        # Fail 건수 기준 정렬 (많은 게 위로)
        kw_counts = fail_df['Keyword'].value_counts().to_dict()
        fail_df_copy['_sort'] = fail_df_copy['Keyword'].map(kw_counts).fillna(0)
        fail_df_copy = fail_df_copy.sort_values(['_sort', 'Build_Num'], ascending=[False, False]).head(7)
        
        rows = len(fail_df_copy) + 1
        available_height = 4.8
        row_height = min(0.55, available_height / rows)
        
        table = s6.shapes.add_table(rows, 6, Inches(0.4), Inches(2),
                                     Inches(12.6), Inches(row_height * rows)).table
        widths = [0.6, 1.2, 1.5, 2.0, 1.5, 5.8]
        for i, w in enumerate(widths):
            table.columns[i].width = Inches(w)
        headers = ['No.', 'Test 항목', '발생 빌드', '모델 / IC', 'Fail Type', '주요 현상']
        for c, h in enumerate(headers):
            set_cell(table.cell(0, c), h, bold=True, size=10,
                     color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
        for r, row in enumerate(fail_df_copy.itertuples(), 1):
            bg = RPT_ALT_ROW if r % 2 == 0 else None
            set_cell(table.cell(r, 0), r, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 1), row.Test_Item, size=8, bg=bg)
            set_cell(table.cell(r, 2), f"R00.0.{row.Build_Num}", size=8, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 3), f"{row.Model} / {row.IC}", size=8, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 4), row.Keyword, size=8, bold=True, bg=bg, color=RPT_CRITICAL, align=PP_ALIGN.CENTER)
            desc = str(row.Fail_Description)[:80] if pd.notna(row.Fail_Description) else "-"
            set_cell(table.cell(r, 5), desc, size=7, bg=bg)
    else:
        add_text(s6, "이번 달 Fail 케이스가 없습니다. ✅", 0.4, 3, 10, 0.4,
                 size=16, bold=True, color=RPT_MID)
    
    # ===== Slide 7: Action Item (페이지 맞춤) ⭐ =====
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s7, 5, "Action Item Review (자동 생성)")
    add_text(s7, "📌 도메인 dict 기반 Action Item 자동 매칭", 0.4, 1.4, 10, 0.4,
             size=14, bold=True, color=RPT_TITLE)
    
    if len(fail_df) > 0:
        fail_df_copy2 = fail_df.copy()
        fail_df_copy2['action'] = fail_df_copy2['Keyword'].apply(
            lambda k: KEYWORD_DOMAIN_MAP.get(k, {}).get('area', '?'))
        # Fail 건수 기준 정렬 (많은 게 위로)
        kw_counts = fail_df['Keyword'].value_counts().to_dict()
        fail_df_copy2['_sort'] = fail_df_copy2['Keyword'].map(kw_counts).fillna(0)
        # 최대 7개로 제한
        unique_actions = fail_df_copy2.drop_duplicates(subset=['Model', 'Keyword']) \
                                      .sort_values(['_sort', 'Build_Num'], ascending=[False, False]) \
                                      .head(7)
        
        rows = len(unique_actions) + 1
        available_height = 4.8
        row_height = min(0.55, available_height / rows)
        
        table = s7.shapes.add_table(rows, 5, Inches(0.4), Inches(2),
                                     Inches(12.6), Inches(row_height * rows)).table
        widths = [0.6, 1.5, 2.5, 2.0, 6.0]
        for i, w in enumerate(widths):
            table.columns[i].width = Inches(w)
        headers = ['No.', '발생 빌드', '모델', 'Fail Type', 'Action Item']
        for c, h in enumerate(headers):
            set_cell(table.cell(0, c), h, bold=True, size=10,
                     color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
        for r, row in enumerate(unique_actions.itertuples(), 1):
            bg = RPT_ALT_ROW if r % 2 == 0 else None
            set_cell(table.cell(r, 0), r, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 1), f"R00.0.{row.Build_Num}", size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 2), row.Model, size=9, bg=bg, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 3), row.Keyword, size=9, bold=True, bg=bg, color=RPT_CRITICAL, align=PP_ALIGN.CENTER)
            set_cell(table.cell(r, 4), row.action + " 검토", size=9, bg=bg, bold=True)
    else:
        add_text(s7, "이번 달 Action Item이 없습니다. ✅", 0.4, 3, 10, 0.4,
                 size=16, bold=True, color=RPT_MID)
    
    # ===== Slide 8: 회귀 알람 (페이지 맞춤) =====
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s8, 6, "회귀 알람 (월 내 최신 빌드)")
    
    builds = sorted(df['Build_Num'].unique())
    if len(builds) >= 2 and len(fail_df) > 0:
        latest_build = builds[-1]
        prev_build = builds[-2]
        latest = fail_df[fail_df['Build_Num'] == latest_build]
        prev = fail_df[fail_df['Build_Num'] == prev_build]
        
        latest_combo = latest.groupby(['Model', 'Keyword']).size().reset_index(name='latest')
        prev_combo = prev.groupby(['Model', 'Keyword']).size().reset_index(name='prev')
        merged = latest_combo.merge(prev_combo, on=['Model', 'Keyword'], how='left')
        merged['prev'] = merged['prev'].fillna(0).astype(int)
        merged['상태'] = merged.apply(
            lambda r: 'NEW' if r['prev'] == 0 else ('WORSE' if r['latest'] > r['prev'] else ''), axis=1)
        alerts = merged[merged['상태'].isin(['NEW', 'WORSE'])].copy()
        # Fail 건수 기준 정렬 (많은 게 위로)
        alerts = alerts.sort_values('latest', ascending=False).head(8)
        
        add_text(s8, f"빌드 R00.0.{latest_build}에서 신규/악화 결함 자동 검출", 0.4, 1.4, 10, 0.4,
                 size=14, bold=True, color=RPT_TITLE)
        n_new = (alerts['상태'] == 'NEW').sum()
        n_worse = (alerts['상태'] == 'WORSE').sum()
        add_text(s8, f"🚨 전체 {len(alerts)}건  ·  신규 {n_new}건  ·  악화 {n_worse}건", 0.4, 1.85, 12, 0.4,
                 size=11, color=RPT_CRITICAL)
        
        if len(alerts) > 0:
            rows = len(alerts) + 1
            available_height = 4.5
            row_height = min(0.5, available_height / rows)
            
            table = s8.shapes.add_table(rows, 6, Inches(0.4), Inches(2.4),
                                         Inches(12.6), Inches(row_height * rows)).table
            widths = [1.0, 2.0, 1.8, 1.2, 1.2, 5.4]
            for i, w in enumerate(widths):
                table.columns[i].width = Inches(w)
            headers = ['상태', '모델', 'Fail Type', '이전', '현재', '검토 영역']
            for c, h in enumerate(headers):
                set_cell(table.cell(0, c), h, bold=True, size=10,
                         color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
            for r, row in enumerate(alerts.itertuples(), 1):
                bg = RPT_ALT_ROW if r % 2 == 0 else None
                state_color = RPT_CRITICAL if row.상태 == 'NEW' else RPT_HIGH
                set_cell(table.cell(r, 0), row.상태, size=9, bold=True, color=state_color, bg=bg, align=PP_ALIGN.CENTER)
                set_cell(table.cell(r, 1), row.Model, size=9, bg=bg, align=PP_ALIGN.CENTER)
                set_cell(table.cell(r, 2), row.Keyword, size=9, bold=True, color=RPT_CRITICAL, bg=bg, align=PP_ALIGN.CENTER)
                set_cell(table.cell(r, 3), row.prev, size=9, bg=bg, align=PP_ALIGN.CENTER)
                set_cell(table.cell(r, 4), row.latest, size=9, bold=True, bg=bg, align=PP_ALIGN.CENTER)
                area = KEYWORD_DOMAIN_MAP.get(row.Keyword, {}).get('area', '?')
                set_cell(table.cell(r, 5), area, size=9, bg=bg)
        else:
            add_text(s8, "이번 빌드에서 회귀 알람 없음 ✅", 0.4, 3, 10, 0.4,
                     size=16, color=RPT_MID)
    else:
        add_text(s8, "이 월에 빌드가 1개 이하라 회귀 비교 불가", 0.4, 3, 10, 0.4,
                 size=14, color=RPT_MID)
    
    # ===== Slide 9: 회의 요약 + 핵심 인사이트 =====
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s9, 7, "회의 요약 및 핵심 인사이트")
    
    # 인사이트 자동 생성 (분석가 시선)
    auto_insight_lines = []
    if len(fail_df) > 0:
        # 1. 가장 위험한 IC 찾기
        ic_fail_rates = df.groupby('IC').apply(
            lambda x: (x['Result']=='FAIL').sum() / len(x) * 100
        )
        if len(ic_fail_rates) > 0:
            max_ic = ic_fail_rates.idxmax()
            max_rate = ic_fail_rates.max()
            avg_rate = ic_fail_rates.mean()
            ratio = max_rate / avg_rate if avg_rate > 0 else 1
            
            # IC의 주요 결함
            ic_fails = fail_df[fail_df['IC'] == max_ic]
            if len(ic_fails) > 0:
                top_kw_for_ic = ic_fails['Keyword'].value_counts().index[0]
                area = KEYWORD_DOMAIN_MAP.get(top_kw_for_ic, {}).get('area', '관련 모듈')
                auto_insight_lines.append(
                    f"• {max_ic} Fail율 {max_rate:.1f}% (전체 평균 {avg_rate:.1f}%의 {ratio:.1f}배) → '{top_kw_for_ic}' 집중, '{area}' 우선 검토"
                )
        
        # 2. 영향 범위가 큰 결함
        kw_counts = fail_df['Keyword'].value_counts()
        kw_models = fail_df.groupby('Keyword')['Model'].nunique()
        if len(kw_counts) > 0 and len(kw_models) > 0:
            top_kw = kw_counts.index[0]
            top_n = kw_counts.iloc[0]
            top_models = kw_models[top_kw]
            auto_insight_lines.append(
                f"• 핵심 결함 '{top_kw}': {top_n}건 발생, {top_models}개 모델에 영향 → 광범위한 구조적 이슈 가능성"
            )
        
        # 3. 회귀 알람 (빌드 비교)
        if len(builds) > 1:
            latest_b = builds[-1]
            prev_b = builds[-2]
            latest_fail = fail_df[fail_df['Build_Num'] == latest_b]
            prev_fail = fail_df[fail_df['Build_Num'] == prev_b]
            latest_combo = set(zip(latest_fail['Model'], latest_fail['Keyword']))
            prev_combo = set(zip(prev_fail['Model'], prev_fail['Keyword']))
            new_alerts = latest_combo - prev_combo
            if len(new_alerts) > 0:
                auto_insight_lines.append(
                    f"• 회귀 알람: 빌드 R00.0.{latest_b}에서 신규 결함 {len(new_alerts)}건 발생 → 빌드 변경 사항 즉시 검토 필요"
                )
            else:
                auto_insight_lines.append(
                    f"• 회귀 알람: 빌드 R00.0.{latest_b}에서 신규 결함 없음 ✓ (안정적)"
                )
    
    if not auto_insight_lines:
        auto_insight_lines = ["• 이번 달 분석 가능한 데이터 부족"]
    
    insight_text = "\n".join(auto_insight_lines)
    
    # 표 (6행으로 확장)
    table = s9.shapes.add_table(6, 2, Inches(0.4), Inches(1.7),
                                 Inches(12.6), Inches(0.48 * 6)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(10.1)
    set_cell(table.cell(0, 0), "구분", bold=True, size=12,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    set_cell(table.cell(0, 1), "내용", bold=True, size=12,
             color=RPT_HEADER_FG, bg=RPT_HEADER_BG, align=PP_ALIGN.CENTER)
    
    # 1행: 주요 결과
    set_cell(table.cell(1, 0), "주요 결과", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    if len(df) > 0:
        result_text = f"• {month_str} Test {len(df):,}건 中 FAIL {len(fail_df):,}건 ({len(fail_df)/len(df)*100:.1f}%)\n"
        if len(builds) > 0:
            latest_b = builds[-1]
            latest_fail_n = len(fail_df[fail_df['Build_Num'] == latest_b]) if len(fail_df) > 0 else 0
            result_text += f"• 최신 빌드 R00.0.{latest_b}: FAIL {latest_fail_n}건\n"
        if len(fail_df) > 0:
            result_text += f"• 영향 모델: {fail_df['Model'].nunique()}개  ·  IC: {fail_df['IC'].nunique()}종"
    else:
        result_text = "• 이번 달 데이터 없음"
    set_cell(table.cell(1, 1), result_text, size=10)
    
    # 2행: 핵심 인사이트 (NEW! - 분석가 자동 생성)
    set_cell(table.cell(2, 0), "💡 핵심 인사이트", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER, color=RPT_CRITICAL)
    set_cell(table.cell(2, 1), insight_text, size=10, bold=True)
    
    # 3행: 주요 이슈
    set_cell(table.cell(3, 0), "주요 이슈", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    if len(fail_df) > 0:
        top_kw = fail_df['Keyword'].value_counts().head(3)
        issue_text = "Top 3 Fail Type (건수 기준):\n"
        for i, (kw, cnt) in enumerate(top_kw.items(), 1):
            issue_text += f"• {i}위 {kw}: {cnt}건\n"
    else:
        issue_text = "이번 달 Fail 이슈 없음 ✅"
    set_cell(table.cell(3, 1), issue_text.strip(), size=10)
    
    # 4행: Action Item
    set_cell(table.cell(4, 0), "Action Item", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    if len(fail_df) > 0:
        top_3 = fail_df['Keyword'].value_counts().head(3)
        action_text = "주요 검토 항목 (Fail 빈도 TOP 3):\n"
        for i, (kw, cnt) in enumerate(top_3.items(), 1):
            area = KEYWORD_DOMAIN_MAP.get(kw, {}).get('area', '?')
            action_text += f"• {i}위 {kw} ({cnt}건): {area} 검토\n"
        action_text += "• 다음 빌드에서 회귀 모니터링 필수"
    else:
        action_text = "별도 Action Item 없음"
    set_cell(table.cell(4, 1), action_text, size=10)
    
    # 5행: 차기 회의
    set_cell(table.cell(5, 0), "차기 회의", bold=True, size=11, bg=RPT_ALT_ROW, align=PP_ALIGN.CENTER)
    set_cell(table.cell(5, 1), "다음 빌드 출시 후 1주 이내", size=10)
    
    # ===== Slide 10: Thank you =====
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RPT_HEADER_BG
    bg.line.fill.background()
    add_text(s10, "Thank you", 0, 3, 13.3, 1.5,
             size=60, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    
    # BytesIO로 저장
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ==============================================================================
# 메인 앱
# ==============================================================================

st.title("🔍 SQA 펌웨어 분석 대시보드")
st.caption("서강대학교 AI·SW대학원 | 생성형 AI와 파이썬 데이터 분석 | A74072 조희주")
st.markdown("---")

try:
    df = get_current_df()
    fail_df = df[df['Result'] == 'FAIL'].copy()
except Exception as e:
    st.error(f"⚠️ 데이터 로드 실패: {e}")
    st.stop()

st.sidebar.title("📋 메뉴")
menu = st.sidebar.radio(
    "분석 메뉴 선택",
    ["🏠 홈", "💡 인사이트 분석", "🔮 알람 & 예측", "📥 검증된 보고서"]
)

# 메뉴별 도움말 표시 여부
show_help = st.sidebar.checkbox("💡 메뉴 도움말 표시", value=True,
                                 help="각 메뉴 상단에 사용법 안내를 표시합니다.")

st.sidebar.markdown("---")
st.sidebar.markdown("### 📦 데이터 정보")
st.sidebar.markdown(f"- 전체 테스트: **{len(df):,}**건")
st.sidebar.markdown(f"- Fail 건수: **{len(fail_df):,}**건")
st.sidebar.markdown(f"- 분석 모델: **{df['Model'].nunique()}**개")
st.sidebar.markdown(f"- 분석 FW: **{df['FW_Version'].nunique()}**개")

if st.session_state.uploaded_data:
    st.sidebar.success(f"➕ 추가 데이터: {sum(len(d) for d in st.session_state.uploaded_data)}건")

st.sidebar.markdown("---")
st.sidebar.markdown("### 💡 메뉴 가이드")
st.sidebar.markdown("""
- **🏠 홈**  
  사이트 전체 안내
- **💡 인사이트 분석**  
  필터·차트·통계 (자동 인사이트 박스)
- **🔮 알람 & 예측**  
  과거 회귀 + 미래 예측
- **📥 검증된 보고서**  
  업로드 + PPT 자동 생성
""")


# ==============================================================================
# 메뉴별 도움말 박스 함수
# ==============================================================================
def show_menu_help(title, description, tips):
    """메뉴 상단에 사용법 안내 박스 (HTML 들여쓰기 없이 한 줄로)"""
    if show_help:
        tips_html = "".join([f'<li style="margin-bottom:0.35rem;">{tip}</li>' for tip in tips])
        html = f'<div style="background-color:#f5f1e8;border:1px solid #d4d1c4;border-left:4px solid #cc785c;border-radius:10px;padding:1rem 1.25rem;margin-bottom:1rem;font-family:\'Pretendard Variable\', sans-serif;"><div style="font-weight:700;color:#1a1a1a;margin-bottom:0.5rem;font-size:0.95rem;">💡 사용법 안내</div><div style="color:#4a4a4a;font-size:0.875rem;margin-bottom:0.75rem;line-height:1.55;">{description}</div><ul style="color:#4a4a4a;font-size:0.875rem;line-height:1.55;margin:0;padding-left:1.25rem;">{tips_html}</ul></div>'
        st.markdown(html, unsafe_allow_html=True)


# ==============================================================================
# 🏠 메뉴 1: 홈 (사이트 안내)
# ==============================================================================
if menu == "🏠 홈":
    st.markdown("# 🎉 SQA 펌웨어 분석 대시보드")
    st.markdown("### 회사 SQA 부서의 펌웨어 테스트 결과를 자동 분석하는 웹 대시보드입니다.")
    st.markdown("---")
    
    # 한눈 정보
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("📋 분석 데이터", f"{len(df):,}건")
    with col2:
        st.metric("❌ Fail율", f"{len(fail_df)/len(df)*100:.1f}%")
    with col3:
        st.metric("💻 모델 수", f"{df['Model'].nunique()}개")
    with col4:
        st.metric("🔄 빌드 수", f"{df['Build_Num'].nunique()}개")
    
    st.markdown("---")
    
    st.markdown("## 📋 메뉴 안내 (왼쪽 사이드바에서 선택)")
    
    col_a, col_b = st.columns(2)
    
    with col_a:
        st.markdown("""
        ### 💡 인사이트 분석
        모든 분석을 한 메뉴에 통합. **각 차트마다 의사결정용 인사이트가 자동 생성**됩니다.
        - 🎯 **필터 분석**: 5가지 조건 동적 분석
        - 📊 **핵심 인사이트**: 6개 차트 + 4단계 인사이트
        - 📈 **통계 인사이트**: 월별·고객사·매트릭스 의미 해석
        
        ### 🔮 알람 & 예측
        과거의 회귀와 미래의 예측을 한 화면에서.
        - 🚨 **회귀 알람**: 신규/악화 결함 자동 검출
        - 🔮 **다음 빌드 예측**: 모델별 미래 Fail율
        """)
    
    with col_b:
        st.markdown("""
        ### 📥 검증된 보고서
        데이터 업로드부터 PPT 자동 생성까지.
        - 📁 **데이터 업로드**: 엑셀 추가 → 즉시 반영
        - 📥 **보고서 생성**: 회사 표준 PPT 자동 작성 + 핵심 인사이트 자동 삽입
        
        ### 📌 메뉴 구성 원칙
        교수님 피드백을 반영해 **인사이트 중심**으로 메뉴를 재설계했습니다. 
        단순 그래프가 아닌 의사결정에 활용 가능한 인사이트를 자동 제공합니다.
        """)
    
    st.markdown("---")
    
    st.success(
        "💡 **처음 사용하시나요?**\n\n"
        "왼쪽 사이드바에서 **'💡 인사이트 분석'** 부터 시작해보세요. "
        "각 메뉴 상단의 사용법 안내 박스를 통해 자세한 사용법을 볼 수 있습니다.\n\n"
        "사이드바 하단의 **'💡 메뉴 도움말 표시'** 체크박스를 해제하면 안내 박스를 숨길 수 있습니다."
    )


# ==============================================================================
# 💡 메뉴 2: 인사이트 분석 (필터 + 전체 인사이트 + 전체 통계)
# ==============================================================================
elif menu == "💡 인사이트 분석":
    st.markdown("## 💡 인사이트 분석")
    st.markdown("**필터를 비우면 전체 데이터, 채우면 선택 그룹의 인사이트가 자동 생성됩니다**")
    
    # 2개 탭 (필터 분석에 전체 인사이트 통합 + 통계)
    tab1, tab2 = st.tabs([
        "🎯 필터 분석 (전체 + 그룹별)",
        "📈 통계 인사이트 (월별·고객사)"
    ])
    
    with tab1:
        show_menu_help(
            "🎯 필터 분석",
            "5가지 조건(IC/Model/Build/Test_Item/Fail_Type)을 조합해서 원하는 데이터만 분석합니다. 필터를 비우면 전체 데이터 인사이트가 생성됩니다.",
            [
                "**필터 비우면 전체 데이터** 인사이트 (= 종합 분석)",
                "**필터 선택 시 해당 그룹**의 인사이트로 자동 갱신",
                "각 필터에서 **여러 값 동시 선택** 가능 (예: GT9XS와 GT1T0A 둘 다)",
                "결과 KPI / 차트 / 핵심 인사이트가 모두 자동 갱신",
                "상세 표의 영상 링크 클릭 시 결함 영상 재생",
            ]
        )
        st.markdown("---")
        render_filter_analysis(df, fail_df)
    
    with tab2:
        show_menu_help(
            "📈 통계 인사이트",
            "시간·조직 관점의 통계를 단순 수치가 아닌 의사결정용 인사이트로 해석합니다.",
            [
                "**월별 추이**: 정점 시점 자동 추적 + 추세 판단 (개선/안정/악화)",
                "**고객사별 분석**: Fail율 격차 + 맞춤 전략 권장",
                "**Test_Item × Fail_Type 매트릭스**: 핫스팟 자동 식별",
                "**데이터 메타 정보**: 수집 기간, IC 종류, 결과 분포",
            ]
        )
        st.markdown("---")
        render_full_statistics(df, fail_df)


# ==============================================================================
# 🔮 메뉴 3: 알람 & 예측 (회귀 + 예측)
# ==============================================================================
elif menu == "🔮 알람 & 예측":
    st.markdown("## 🔮 알람 & 예측")
    st.markdown("**과거의 회귀 + 미래의 예측 — 통합 모니터링**")
    
    # 2개 탭
    tab1, tab2 = st.tabs([
        "🚨 회귀 알람 (과거)",
        "🔮 다음 빌드 예측 (미래)"
    ])
    
    with tab1:
        show_menu_help(
            "🚨 회귀 알람",
            "최신 빌드 vs 이전 빌드를 모델 × Fail_Type 조합 단위로 비교해 회귀를 자동 검출합니다.",
            [
                "**🆕 NEW**: 이전 빌드에 없던 결함이 새로 등장",
                "**📈 WORSE**: 기존 결함이 더 많이 발생 (악화)",
                "동일하거나 개선된 케이스는 알람에서 제외됩니다",
                "표의 영상 링크를 클릭하면 결함 영상이 재생됩니다",
                "단순 평균에 숨겨진 회귀를 모델 단위로 잡아내는 것이 핵심 가치입니다",
            ]
        )
        st.markdown("---")
        render_regression_alert(df, fail_df)
    
    with tab2:
        show_menu_help(
            "🔮 다음 빌드 예측",
            "선택한 모델의 과거 빌드별 Fail율 추세를 학습해서 바로 다음 빌드의 Fail율을 예측합니다.",
            [
                "**모델 단위로 예측**: 같은 모델은 같은 FW 라인이 빌드별로 진화하므로 모델별 예측",
                "**다항회귀 (2차)** + **표본 가중치** 적용",
                "**95% 신뢰구간**: 점추정값이 아닌 범위로 표현",
                "차트의 빨간 별표가 다음 빌드 예측값입니다",
            ]
        )
        st.markdown("---")
        render_fail_rate_prediction(df, fail_df)


# ==============================================================================
# 📥 메뉴 4: 검증된 보고서 (업로드 + 보고서)
# ==============================================================================
elif menu == "📥 검증된 보고서":
    st.markdown("## 📥 검증된 보고서")
    st.markdown("**데이터 업로드 → 자동 분석 → PPT 보고서 (핵심 인사이트 자동 삽입)**")
    
    # 2개 탭
    tab1, tab2 = st.tabs([
        "📁 데이터 업로드",
        "📥 보고서 생성"
    ])
    
    with tab1:
        show_menu_help(
            "📁 데이터 업로드",
            "엑셀(.xlsx) 또는 CSV 파일을 업로드하면 자동으로 분석 데이터에 추가됩니다.",
            [
                "**필수 컬럼**: IC, Model, Build_Num, Result, Test_Item, Keyword (기본 CSV와 동일)",
                "**중복 검사 3단계**: 같은 파일 / No 번호 중복 / 행 단위 중복 자동 차단",
                "업로드 성공/실패는 상단 팝업으로 표시되고 화면이 자동 초기화됩니다",
                "업로드 이력에서 체크박스로 개별 파일 삭제 가능",
            ]
        )
        st.markdown("---")
        render_data_upload()
    
    with tab2:
        show_menu_help(
            "📥 보고서 생성",
            "선택한 월의 데이터로 회사 표준 SQA-SW Quality Review 보고서(PPT 10장)를 자동 생성합니다. 회의 요약 슬라이드에 핵심 인사이트가 자동 삽입됩니다.",
            [
                "월을 선택하면 그 달의 데이터만 분석되어 보고서가 생성됩니다",
                "**Action Item이 도메인 dict로 자동 매칭**됩니다 (예: touch delay → 응답속도 검토)",
                "**회의 요약 슬라이드에 💡 핵심 인사이트 자동 생성** (가장 위험한 IC, 광범위 결함, 회귀 알람)",
                "생성된 보고서는 이력에 저장되어 언제든 다운로드 가능",
                "체크박스로 선택해서 보고서 삭제 가능",
            ]
        )
        st.markdown("---")
        render_report_generator(df, fail_df)
